#!/usr/bin/env python3
"""APPEND one slide on the MATSim assignment + skim-feedback loop to the (user-edited)
deck, WITHOUT regenerating anything else. Opens VAE_Architecture_and_Validation.pptx,
adds a single plain-style slide at the end, saves in place.

IMPORTANT: close PowerPoint first, or the save will clobber it.
Optional: python append_matsim_slide.py <in_out.pptx>   (dry-run on a copy)
"""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
DECK = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "VAE_Architecture_and_Validation.pptx"

BLACK = RGBColor(0, 0, 0); GREY = RGBColor(0x55, 0x55, 0x55); LINE = RGBColor(0x33, 0x33, 0x33)
prs = Presentation(str(DECK))
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
MARGIN = Inches(0.6)


def _tb(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
    return tf


def title(s, text, sub=None, size=24):
    tf = _tb(s, MARGIN, Inches(0.22), W - 2 * MARGIN, Inches(1.0))
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    rule_y = Inches(0.78)
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(13); rr.font.italic = True; rr.font.name = "Times New Roman"; rr.font.color.rgb = GREY
        rule_y = Inches(1.12)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, rule_y, W - 2 * MARGIN, Pt(1.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    return rule_y


def bullets(s, items, top=Inches(1.5), size=15):
    tf = _tb(s, MARGIN, top, W - 2 * MARGIN, H - top - Inches(0.4))
    for i, it in enumerate(items):
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(5)
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size - 2 * lvl); r.font.name = "Times New Roman"
        r.font.color.rgb = BLACK if lvl == 0 else GREY


def block(s, l, t, w, h, text, fontsize=12, bold=False):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(255, 255, 255)
    b.line.color.rgb = BLACK; b.line.width = Pt(1); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(fontsize); r.font.bold = bold
    r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    return b


def arrow(s, l, t, w, kind=MSO_SHAPE.RIGHT_ARROW, h=Inches(0.28)):
    a = s.shapes.add_shape(kind, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
    a.line.fill.background(); a.shadow.inherit = False


def lbl(s, l, t, w, text, italic=True, size=11, align=PP_ALIGN.CENTER):
    tf = _tb(s, l, t, w, Inches(0.32)); p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.italic = italic
    r.font.name = "Times New Roman"; r.font.color.rgb = GREY


# ── the MATSim feedback-loop slide ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title(s, "MATSim — Traffic Assignment & the Skim-Feedback Loop",
      "route assignment returns congested skims that re-drive MITO mode & destination choice → policy-sensitive mode shift")

# two-box clockwise loop: MITO ⇄ MATSim
bh = Inches(1.35)
bL = block(s, Inches(0.9), Inches(1.75), Inches(4.8), bh,
           "MITO — tour-based demand\nmode & destination choice → one chained daily plan per person", 12, True)
bR = block(s, Inches(7.65), Inches(1.75), Inches(4.8), bh,
           "MATSim — dynamic traffic assignment\ncar mobsim + mapped PT; agents re-route to a network equilibrium", 12, True)
# top arrow → (plans)
arrow(s, Inches(5.75), Inches(1.95), Inches(1.9), MSO_SHAPE.RIGHT_ARROW)
lbl(s, Inches(5.55), Inches(1.62), Inches(2.3), "population plans file")
# bottom arrow ← (skims)
arrow(s, Inches(5.75), Inches(2.85), Inches(1.9), MSO_SHAPE.LEFT_ARROW)
lbl(s, Inches(5.35), Inches(3.16), Inches(2.7), "updated congested skims\n(zone-to-zone times & costs)")

bullets(s, [
 ("The tour-based model writes one chained daily plan per person; MATSim assigns those plans on a PT-mapped "
  "multimodal network (OSM road + MTA GTFS via pt2matsim) and iterates agents to a route/departure-time equilibrium.", 0),
 ("MATSim returns congested zone-to-zone travel-time and cost skims; these replace the fixed input skims and are fed "
  "back into MITO's mode and destination choice, which re-run — the MITO ⇄ MATSim loop iterates to convergence.", 0),
 ("Assignment is validated against 2017 MDOT SHA AADT count stations (target GEH < 5) for the 6-county Baltimore region.", 0),
 ("Policy sensitivity — the point of the loop: mode choice reads level-of-service from the skims, so a policy that "
  "changes the network (e.g. an I-695 toll, or added congestion) updates the skims and makes mode choice re-respond, "
  "producing genuine mode shift. Without the loop the skims are fixed and the model cannot react — the loop is what "
  "lets the VOT-based cost coefficient actually bite.", 0),
], top=Inches(3.75), size=14)

prs.save(str(DECK))
print(f"appended MATSim slide → {DECK}  (now {len(prs.slides._sldIdLst)} slides)")
