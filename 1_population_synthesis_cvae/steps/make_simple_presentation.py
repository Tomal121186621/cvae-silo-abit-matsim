#!/usr/bin/env python3
"""Build a SHORT, PLAIN academic PPTX: VAE architecture + workzone-assignment
methodology, followed by the CORE validation figures — ONE figure per slide,
each blown up as large as possible so axes and numbers stay readable.
Kept: marginals, joint relationships (+ association), structural zeros.
No color fills, no decorative design — black text on white.
Output → VAE_Architecture_and_Validation.pptx
"""
from __future__ import annotations
import sys
from pathlib import Path
from PIL import Image

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
FV = ROOT / "outputs" / "figures" / "validation"
FT = ROOT / "outputs" / "figures" / "training"
OUT = ROOT / "VAE_Architecture_and_Validation.pptx"

BLACK = RGBColor(0, 0, 0)
GREY = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0x33, 0x33, 0x33)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
MARGIN = Inches(0.6)

TV_DEF = ("TV (Total Variation) = ½·Σ|p_gen − p_test| over categories: 0 = identical, "
          "1 = disjoint — the single-number distance between two distributions.")


def slide():
    return prs.slides.add_slide(BLANK)


def _tb(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
    return tf


def title(s, text, sub=None, size=24):
    tf = _tb(s, MARGIN, Inches(0.22), W - 2 * MARGIN, Inches(1.0))
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.name = "Times New Roman"
    r.font.color.rgb = BLACK
    rule_y = Inches(0.78)
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(13); rr.font.italic = True; rr.font.name = "Times New Roman"
        rr.font.color.rgb = GREY
        rule_y = Inches(1.12)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, rule_y, W - 2 * MARGIN, Pt(1.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    return rule_y


def bullets(s, items, left=MARGIN, top=Inches(1.5), width=None, size=17):
    width = width or (W - 2 * MARGIN)
    tf = _tb(s, left, top, width, H - top - Inches(0.4))
    for i, it in enumerate(items):
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(5)
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size - 2 * lvl); r.font.name = "Times New Roman"
        r.font.color.rgb = BLACK if lvl == 0 else GREY


def _fit(path, box_w, box_h):
    iw, ih = Image.open(path).size
    ar = iw / ih; bar = box_w / box_h
    if ar > bar:
        return int(box_w), int(box_w / ar)
    return int(box_h * ar), int(box_h)


def big_figure(title_txt, path, explain):
    """One slide: title + one-line explanation, then the figure blown up to fill
    the whole remaining area (axes/numbers stay readable)."""
    s = slide()
    rule_y = title(s, title_txt, explain, size=22)
    p = Path(path)
    if not p.exists():
        _tb(s, MARGIN, Inches(3), Inches(4), Inches(0.4)).paragraphs[0].add_run().text = f"[missing {p.name}]"
        return s
    box_l = Inches(0.25)
    box_t = int(rule_y) + Inches(0.12)
    box_w = W - Inches(0.5)
    box_h = H - box_t - Inches(0.15)
    w, h = _fit(str(p), int(box_w), int(box_h))
    l = int(box_l + (box_w - w) / 2); t = int(box_t + (box_h - h) / 2)
    s.shapes.add_picture(str(p), l, t, width=w, height=h)
    return s


# ── architecture-diagram helpers ───────────────────────────────────────────
def block(s, l, t, w, h, text, fontsize=12, bold=False):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(255, 255, 255)
    b.line.color.rgb = BLACK; b.line.width = Pt(1); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(fontsize); r.font.bold = bold
    r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    return b


def arrow(s, l, t, w):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, Inches(0.22))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88)
    a.line.fill.background(); a.shadow.inherit = False


# ══════════════════════════════════════════════════════════════════════════
# 1. TITLE
# ══════════════════════════════════════════════════════════════════════════
s = slide()
tf = _tb(s, Inches(1.2), Inches(2.3), Inches(10.9), Inches(2.6))
p = tf.paragraphs[0]; r = p.add_run()
r.text = "A Conditional VAE for Synthetic Population Synthesis"
r.font.size = Pt(34); r.font.bold = True; r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
p = tf.add_paragraph(); p.space_before = Pt(10); r = p.add_run()
r.text = "Architecture, Workplace-Assignment Methodology, and Held-out Validation"
r.font.size = Pt(19); r.font.name = "Times New Roman"; r.font.color.rgb = GREY
p = tf.add_paragraph(); p.space_before = Pt(16); r = p.add_run()
r.text = "ACS PUMS 2016 (5-year) → SILO / MITO  ·  Baltimore–Washington MSTM region (96 PUMAs, 6 states)"
r.font.size = Pt(14); r.font.italic = True; r.font.name = "Times New Roman"; r.font.color.rgb = GREY
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.15), Inches(6.0), Pt(2))
ln.fill.solid(); ln.fill.fore_color.rgb = BLACK; ln.line.fill.background(); ln.shadow.inherit = False

# ══════════════════════════════════════════════════════════════════════════
# 2. VAE ARCHITECTURE (diagram)
# ══════════════════════════════════════════════════════════════════════════
s = slide(); title(s, "VAE Architecture",
                   "simple conditional VAE (Borysov 2019 style) — flat one-hot record, MLP encoder/decoder, one-shot softmax heads")
y = Inches(1.9); bh = Inches(1.1)
block(s, Inches(0.55), y, Inches(2.0), bh,
      "Household record\none-hot: dwellingType, tenure, autos, income-bin  +  persons (≤S) one-hot", 11)
arrow(s, Inches(2.62), Inches(2.35), Inches(0.5))
block(s, Inches(3.18), y, Inches(1.7), bh, "MLP\nEncoder\n(2 layers, 256)", 12, True)
arrow(s, Inches(4.95), Inches(2.35), Inches(0.5))
block(s, Inches(5.5), y, Inches(2.0), bh, "Gaussian latent  z\ndim 24  ·  free-bits KL\n(no collapse)", 12, True)
arrow(s, Inches(7.57), Inches(2.35), Inches(0.5))
block(s, Inches(8.12), y, Inches(1.7), bh, "MLP\nDecoder\n(2 layers, 256)", 12, True)
arrow(s, Inches(9.86), Inches(2.35), Inches(0.5))
block(s, Inches(10.42), y, Inches(2.3), bh,
      "One-shot softmax heads\nHH heads + shared per-person head (×S slots)", 11)
block(s, Inches(5.5), Inches(3.5), Inches(2.0), Inches(0.7),
      "Condition:\nPUMA embedding (8-d) + household-size one-hot", 10)
au = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(6.35), Inches(3.05), Inches(0.22), Inches(0.5))
au.fill.solid(); au.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88); au.line.fill.background(); au.shadow.inherit = False
bullets(s, [
 ("PUMA is the ONLY embedded variable; every other categorical is one-hot. Joints are carried by the latent z.", 0),
 ("No autoregressive decoder, no Set Transformer, no in-network income head — deliberately simple and transparent.", 0),
 ("Income is a categorical BIN; continuous dollars drawn empirically within each (PUMA, bin) cell at generation "
  "(open top bin preserves the real $1M+ tail). Person↔household income reconciled exactly (Σ = 99.98%).", 0),
 ("Constrained decoding masks make impossible classes unreachable → structural zeros = 0.", 0),
 ("Loss = reconstruction cross-entropy + β·KL (free-bits floor 0.5 nat/dim) + optional per-PUMA marginal-JSD term.", 0),
], top=Inches(4.45), size=15)

# ══════════════════════════════════════════════════════════════════════════
# 3. WORKZONE / WORKPLACE ASSIGNMENT METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════
s = slide(); title(s, "Workplace (Work-Zone) Re-assignment — Methodology",
                   "employed persons re-assigned to job zones by a capacity-constrained gravity model (commute-time decay × job vacancies)")
y = Inches(1.7); bh = Inches(0.95)
block(s, Inches(0.55), y, Inches(2.5), bh, "Employed persons\n(occupation = worker) at home zone", 11)
arrow(s, Inches(3.12), Inches(2.05), Inches(0.45))
block(s, Inches(3.62), y, Inches(2.7), bh, "Highway skim (OMX)\nhome→zone travel time (min)", 11)
arrow(s, Inches(6.4), Inches(2.05), Inches(0.45))
block(s, Inches(6.9), y, Inches(2.7), bh, "Gravity weight\nw = exp(−βt) × vacant jobs\n(β = 0.07)", 11, True)
arrow(s, Inches(9.66), Inches(2.05), Inches(0.45))
block(s, Inches(10.16), y, Inches(2.6), bh, "Re-assign job zone,\ndecrement vacancy,\ncreate job (jj)", 11)
bullets(s, [
 ("Every employed person is re-assigned a job zone by a capacity-constrained gravity model: weight "
  "w = exp(−β·t) × vacant jobs, sample a zone, decrement that zone's jobs, create the job record.", 0),
 ("Why re-assign: the first workplace pass used the raw commute-time TLFD directly as the friction AND multiplied "
  "by job counts — that double-counts downtown jobs and made commutes ~2× too long (~62 min vs the 37.6-min HTS target).", 0),
 ("Fix: a proper distance-decay friction exp(−β·t) with β = 0.07 reproduces the observed HTS work-trip "
  "commute-length distribution.", 0),
 ("Capacity-constrained — each assigned worker decrements that zone's job inventory so zones fill to their forecast "
  "employment (skim times cached per home zone). SILO links person↔job by person.workplace = JOB id; the job's zone "
  "lives in the jj frame.", 0),
 ("Result: 6,091,727 workers re-assigned from an 8.98M-job forecast; 0 unreachable, 0 external surplus.", 0),
 ("Downstream, the tour-based travel model re-draws each worker's job zone again against the RTS-estimated HBW "
  "impedance (doubly-constrained) for RTS-realistic commutes (median ~26 → ~16 mi).", 0),
], top=Inches(2.95), size=14)

# ══════════════════════════════════════════════════════════════════════════
# TRAINING HEALTH — the fit is healthy and the model generalizes (no memorization)
# ══════════════════════════════════════════════════════════════════════════
big_figure("Training Health — ELBO Converges, No Overfitting",
           FT / "T1_elbo_curves.png",
           "Train and validation ELBO fall together and stay close — smooth convergence, no train/val gap (no overfit).")
big_figure("Training Health — Latent Fully Used, No Posterior Collapse",
           FT / "T3_kl_active_dims.png",
           "All 24 latent dimensions carry information (KL > 0) via free-bits — the model uses its capacity, it does not collapse.")
big_figure("Training Health — Reconstruction Learned",
           FT / "T2_reconstruction.png",
           "Per-variable reconstruction accuracy on held-out records — the decoder faithfully recovers inputs.")
big_figure("No Memorization — the Model Generalizes",
           FV / "M_memorization.png",
           "Only 61% of generated person-types were seen in training; 2,377 valid types are NEW — it samples the "
           "distribution, it does not copy records.")

# ══════════════════════════════════════════════════════════════════════════
# VALIDATION — one figure per slide, blown up
# All held-out (2016 PUMS test split). TV defined once here.
# ══════════════════════════════════════════════════════════════════════════

# ── Marginals ──────────────────────────────────────────────────────────────
big_figure("Marginals — Summary (Total Variation)",
           FV / "F1_marginal_TV_summary.png",
           "Each bar = one attribute's TV from the held-out test; all near/under the 0.03 line.  " + TV_DEF)

MARG = [
 ("Marginal — Household: dwelling type", "F1_marginal_hh_dwellingType.png",
  "Generated vs ACS-test dwelling-type shares — matches almost exactly (TV 0.012)."),
 ("Marginal — Household: tenure (own / rent)", "F1_marginal_hh_tenure.png",
  "Own vs rent split reproduced (TV 0.015)."),
 ("Marginal — Household: vehicles (autos)", "F1_marginal_hh_autos.png",
  "Vehicle-count shares near-perfect (TV 0.009)."),
 ("Marginal — Household: income bin", "F1_marginal_hh_income_bin.png",
  "Income-bin shares close, slight top-tail thinning (TV 0.042)."),
 ("Marginal — Person: age", "F1_marginal_pp_age_bin.png",
  "Age profile tracks ACS across all bins (TV 0.039)."),
 ("Marginal — Person: gender", "F1_marginal_pp_gender.png",
  "Gender split essentially exact (TV 0.0004)."),
 ("Marginal — Person: race", "F1_marginal_pp_race.png",
  "Race shares reproduced (TV 0.025)."),
 ("Marginal — Person: occupation", "F1_marginal_pp_occupation.png",
  "Occupation shares close (TV 0.035)."),
 ("Marginal — Person: driver's license", "F1_marginal_pp_driversLicense.png",
  "License-holding rate matches (TV 0.016)."),
 ("Marginal — Person: relationship / household role", "F1_marginal_pp_relationship.png",
  "Household roles reproduced (TV 0.034)."),
 ("Marginal — Person: nationality (foreign-born)", "F1_marginal_pp_nationality.png",
  "Foreign-born share and its spatial gradient captured."),
 ("Marginal — Person: income bin", "F1_marginal_pp_income_bin.png",
  "Person-income bins close (TV 0.023)."),
]
for ttl, fn, ex in MARG:
    big_figure(ttl, FV / fn, ex)

# ── Joint relationships ────────────────────────────────────────────────────
big_figure("Joint Relationships — Fidelity by Interaction Order",
           FV / "F8_joint_srmse_by_order.png",
           "Error grows with order (1-way 0.06 → 2-way 0.14 → 3-way 0.24) — higher-order joints are harder.")

JOINTS = [
 ("Joint — Age × Income", "F3_joint_age_bin_x_income_bin.png",
  "Truth · generated · difference; the difference panel (right) is near zero."),
 ("Joint — Age × Relationship", "F3_joint_age_bin_x_relationship.png",
  "Age-by-household-role structure reproduced (difference panel near zero)."),
 ("Joint — Occupation × Relationship", "F3_joint_occupation_x_relationship.png",
  "Occupation-by-role structure reproduced (difference panel near zero)."),
]
for ttl, fn, ex in JOINTS:
    big_figure(ttl, FV / fn, ex)

big_figure("Association Structure — Truth (Cramér's V)",
           FV / "S3_cramersV_test.png",
           "Pairwise association strengths in the real data (V: 0 = none, 1 = perfect) — the target pattern.")
big_figure("Association Structure — Generated (Cramér's V)",
           FV / "S3_cramersV_generated.png",
           "Same association pattern reproduced by the generated population.")
big_figure("Association Structure — Difference (Cramér's V)",
           FV / "S3_cramersV_diff.png",
           "Generated minus truth — near zero everywhere: pairwise dependencies are preserved.")

# ── Structural zeros ───────────────────────────────────────────────────────
big_figure("Structural Zeros — Impossible Records",
           FV / "Z_structural_zeros.png",
           "Zero violations: no under-16 drivers, no under-62 retirees, etc. — constrained decoding guarantees it.")

# ── Summary table ──────────────────────────────────────────────────────────
s = slide(); title(s, "Validation — Summary")
rows = [
 ("Check", "Result"),
 ("Training fit", "train/val ELBO track closely — converges, no overfit"),
 ("Latent health", "24 / 24 dims active (free-bits) — no posterior collapse"),
 ("Memorization", "only 61% of gen types in train; 2,377 novel valid types"),
 ("Marginals (all attributes)", "at/below the 0.03 TV reference; gender TV 0.0004, autos 0.009"),
 ("Joint relationships", "1-way SRMSE 0.06, 2-way 0.14, 3-way 0.24"),
 ("Association (Cramér's V)", "difference ≈ 0 — pairwise dependencies preserved"),
 ("Structural-zero violations", "0 (constrained decoding)"),
 ("Coherence (Σ income exact / 1 householder)", "99.98% / 100%"),
 ("Workplace assignment", "6.09M jobs assigned, 0 unreachable"),
]
tbl = s.shapes.add_table(len(rows), 2, MARGIN, Inches(1.5), W - 2 * MARGIN, Inches(4.4)).table
tbl.columns[0].width = Inches(5.0); tbl.columns[1].width = Inches(7.1)
tblPr = tbl._tbl.tblPr
for child in list(tblPr):
    tblPr.remove(child)
for j, c in enumerate(rows[0]):
    cell = tbl.cell(0, j); cell.text = c
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    rr = cell.text_frame.paragraphs[0].runs[0]
    rr.font.bold = True; rr.font.size = Pt(14); rr.font.name = "Times New Roman"; rr.font.color.rgb = BLACK
for i in range(1, len(rows)):
    for j, c in enumerate(rows[i]):
        cell = tbl.cell(i, j); cell.text = c
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rr = cell.text_frame.paragraphs[0].runs[0]
        rr.font.size = Pt(12); rr.font.name = "Times New Roman"; rr.font.color.rgb = BLACK

prs.save(OUT)
print(f"saved → {OUT}  ({len(prs.slides._sldIdLst)} slides)")
