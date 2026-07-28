#!/usr/bin/env python3
"""APPEND a Tour-Based MITO section to the (user-edited) deck WITHOUT regenerating it.
Opens VAE_Architecture_and_Validation.pptx, appends MITO methodology (plain diagrams)
+ validation figures (one clean figure per slide, VAE style), saves in place.

IMPORTANT: close PowerPoint first (deck must not be open) or your save will clobber these.
Optional: python append_mito_slides.py <in_out.pptx>   (dry-run on a copy)
"""
from __future__ import annotations
import sys
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
DECK = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "VAE_Architecture_and_Validation.pptx"
MITO = ROOT.parent / "Tour Based MITO"
MFIG = MITO / "validation" / "figures"
MAPP = MFIG / "applied"

BLACK = RGBColor(0, 0, 0); GREY = RGBColor(0x55, 0x55, 0x55); LINE = RGBColor(0x33, 0x33, 0x33)

prs = Presentation(str(DECK))
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
MARGIN = Inches(0.6)


def slide():
    return prs.slides.add_slide(BLANK)


def _tb(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
    return tf


def title(s, text, sub=None, size=24):
    tf = _tb(s, MARGIN, Inches(0.22), W - 2 * MARGIN, Inches(1.0))
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    rule_y = Inches(0.78)
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(13); rr.font.italic = True; rr.font.name = "Times New Roman"; rr.font.color.rgb = GREY
        rule_y = Inches(1.12)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, rule_y, W - 2 * MARGIN, Pt(1.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    return rule_y


def bullets(s, items, left=MARGIN, top=Inches(1.5), width=None, size=15):
    width = width or (W - 2 * MARGIN)
    tf = _tb(s, left, top, width, H - top - Inches(0.4))
    for i, it in enumerate(items):
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(5)
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size - 2 * lvl); r.font.name = "Times New Roman"
        r.font.color.rgb = BLACK if lvl == 0 else GREY


def block(s, l, t, w, h, text, fontsize=12, bold=False):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(255, 255, 255)
    b.line.color.rgb = BLACK; b.line.width = Pt(1); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(fontsize); r.font.bold = bold
    r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    return b


def arrow(s, l, t, w, kind=MSO_SHAPE.RIGHT_ARROW, h=Inches(0.2)):
    a = s.shapes.add_shape(kind, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88)
    a.line.fill.background(); a.shadow.inherit = False


def flow_row(s, labels, y, box_w, box_h, gap, x0=Inches(0.35), bold=False, fs=11):
    """Place a horizontal row of boxes separated by right-arrows; return end x."""
    x = int(x0); step = int(box_w) + int(gap)
    for i, lab in enumerate(labels):
        block(s, x, int(y), int(box_w), int(box_h), lab, fs, bold)
        if i < len(labels) - 1:
            arrow(s, x + int(box_w) + int((int(gap) - int(Inches(0.28))) / 2),
                  int(y) + int(box_h) // 2 - int(Inches(0.1)), Inches(0.28))
        x += step
    return x


def big_figure(title_txt, path, explain):
    s = slide(); rule_y = title(s, title_txt, explain, size=22)
    p = Path(path)
    if not p.exists():
        _tb(s, MARGIN, Inches(3), Inches(6), Inches(0.4)).paragraphs[0].add_run().text = f"[missing {p.name}]"
        return s
    iw, ih = Image.open(str(p)).size; ar = iw / ih
    box_l = Inches(0.25); box_t = int(rule_y) + Inches(0.12)
    box_w = W - Inches(0.5); box_h = H - box_t - Inches(0.15)
    if ar > box_w / box_h:
        w, h = int(box_w), int(box_w / ar)
    else:
        w, h = int(box_h * ar), int(box_h)
    l = int(box_l + (box_w - w) / 2); t = int(box_t + (box_h - h) / 2)
    s.shapes.add_picture(str(p), l, t, width=w, height=h)
    return s


def section(text, subtitle=None):
    s = slide()
    tf = _tb(s, MARGIN, Inches(2.8), W - 2 * MARGIN, Inches(1.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(4.0), Inches(5.3), Pt(1.4))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    if subtitle:
        tf2 = _tb(s, MARGIN, Inches(4.2), W - 2 * MARGIN, Inches(1.2))
        p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = subtitle
        r.font.size = Pt(14); r.font.italic = True; r.font.name = "Times New Roman"; r.font.color.rgb = GREY
    return s


# ══════════════════════════════════════════════════════════════════════════
# MITO — divider
# ══════════════════════════════════════════════════════════════════════════
section("Tour-Based MITO — Travel-Demand Model",
        "home-anchored tours estimated on the RTS 2017–18 travel survey; applied on the SILO/VAE population → MATSim")

# ── 1. Overview / why tour-based ───────────────────────────────────────────
s = slide(); title(s, "Tour-Based MITO — Overview",
                   "the unit of demand is the home-anchored TOUR (home → stops → primary activity → stops → home)")
y = Inches(1.75); bh = Inches(1.05)
block(s, Inches(0.55), y, Inches(2.7), bh, "SILO / VAE\nsynthetic population\n(persons · households)", 11)
arrow(s, Inches(3.32), Inches(2.15), Inches(0.5))
block(s, Inches(3.88), y, Inches(3.0), bh, "Tour-based demand model\n(9 estimated components)", 11, True)
arrow(s, Inches(6.94), Inches(2.15), Inches(0.5))
block(s, Inches(7.5), y, Inches(2.5), bh, "One chained daily plan\nper person", 11)
arrow(s, Inches(10.06), Inches(2.15), Inches(0.5))
block(s, Inches(10.62), y, Inches(2.15), bh, "MATSim\nassignment", 11)
bullets(s, [
 ("Replaces trip-based MITO, which over-generated travel (~6.3 legs/person vs the survey's 3.9) because home-based "
  "rates were calibrated to one-way legs but materialized as 2-leg tours.", 0),
 ("Fix: calibrate to TOURS per person (grounded in RTS) → legs fall out at the right count; one chained plan per "
  "person by construction — exactly what MATSim needs.", 0),
 ("Estimated on the Regional Travel Survey 2017–18 (163,290 trips / 41,914 persons / 21,788 households, "
  "Baltimore–Washington).", 0),
 ("A single tour-level mode choice with a value-of-time cost coefficient is the natural home for congestion-pricing "
  "elasticity (the I-695 study).", 0),
], top=Inches(3.15), size=14)

# ── 2. The tour pipeline (9 components) ────────────────────────────────────
s = slide(); title(s, "The Tour Pipeline — 9 Estimated Components",
                   "each person → daily pattern → tours → for every tour: destination, mode, stops, schedule")
row1 = ["0 · Daily activity\npattern (MNL)", "1 · Tour frequency\nby purpose", "2 · Tour destination\n(gravity MNL)",
        "3 · Tour mode choice\n(VOT-based MNL)", "4 · Stop frequency\n(ordered logit)"]
row2 = ["5 · Stop location\n(detour MNL)", "6 · Time of day\n(schedule)", "7 · Travel-time budget\n(Weibull AFT)",
        "8 · Stop purpose", "→ Chained plans\n(MATSim)"]
flow_row(s, row1, Inches(1.55), Inches(2.32), Inches(0.95), Inches(0.34), fs=10.5)
flow_row(s, row2, Inches(3.05), Inches(2.32), Inches(0.95), Inches(0.34), fs=10.5)
bullets(s, [
 ("Estimation: maximum likelihood — Apollo for the discrete-choice models, survival::survreg (Weibull AFT) for the "
  "travel-time budget; survey-weighted with robust (sandwich) standard errors.", 0),
 ("Validation split: household-level 80/20 train/hold-out (members never straddle the split); 27 estimations, all converged.", 0),
 ("Apply-time calibration is standard and preserves the estimated behavioural parameters: tour/stop rates to RTS, a "
  "per-purpose gravity length multiplier, and mode ASCs — the VOT cost coefficient (the pricing elasticity) is left untouched.", 0),
], top=Inches(4.35), size=14)

# ── 3. Behavioural foundations ─────────────────────────────────────────────
s = slide(); title(s, "Behavioural Foundations",
                   "the modelling ideas that make the tour model realistic and policy-usable")
bullets(s, [
 ("Tour as the unit of generation — one home-anchored chain, so trip counts are correct by construction.", 0),
 ("One mode per tour with a value-of-time cost coefficient — income-dependent VOT, Maryland-calibrated to ~$22/hr at "
  "the regional median, with the correct equity gradient (low income → lower VOT → more toll-deterred).", 0),
 ("Destination choice = distance decay (b_time) + attraction size — chosen destinations average ~31 min vs ~84 min "
  "for random zones; work tours decay slowest (people travel farther to work).", 0),
 ("Travel-time budget (Zahavi regularity): a roughly constant ~80 min/person/day daily budget (Weibull AFT) that "
  "constrains destination + stop choices — the grounded replacement for a rubber-band heuristic.", 0),
 ("Two-peak time-of-day: the schedule reproduces the AM/PM commute peaks.", 0),
 ("Workplace re-assignment: SILO commutes ran ~1.7× too long (median ~26 mi vs RTS ~16), so each worker's job zone is "
  "re-drawn against the RTS work-trip impedance (doubly-constrained, job totals honoured).", 0),
], size=15)

# ══════════════════════════════════════════════════════════════════════════
# MITO — validation figures (one clean figure per slide)
# ══════════════════════════════════════════════════════════════════════════
# two estimation-diagnostic figures, then the applied (system) validation set vs the RTS survey
FIGS = [
 (MFIG / "fig9_goodness_of_fit.png", "Estimation Fit — McFadden ρ² by Component",
  "Bars = each component's McFadden ρ² (0 = no better than chance; higher = better). All 9 components fit well "
  "(ρ² ≈ 0.27–0.81); every one of the 27 estimations converged."),
 (MFIG / "fig3_oos_first_pref_recovery.png", "Out-of-Sample — First-Preference Recovery vs Chance",
  "Held-out 20% of households: how often the model's top-choice matches the observed choice, vs the chance baseline. "
  "Every component beats chance — DAP 0.79 (chance 0.50); mode HBW/HBS/HBO 0.65/0.66/0.56; destination up to 19× chance."),
 (MAPP / "val1_purpose_share.png", "Applied vs Survey — Tour-Purpose Share",
  "Share of tours by purpose (HBW/HBS/HBO): the applied model on the SILO/VAE population matches the RTS survey within ~1 pp."),
 (MAPP / "val2_mode_by_purpose.png", "Applied vs Survey — Mode Share by Purpose",
  "Mode shares for each purpose (HBW/HBS/HBO); the applied bars track the survey across auto, transit, walk and bike."),
 (MAPP / "val3_time_of_day.png", "Applied vs Survey — Time of Day",
  "Outbound and return departure-hour distributions; the applied schedule reproduces the AM and PM commute peaks."),
 (MAPP / "val4_trip_length_by_purpose.png", "Applied vs Survey — Trip Length by Purpose",
  "Primary-distance distributions by purpose; shapes match, with applied commutes running slightly shorter than the survey (≈16 vs 20 mi for work)."),
 (MAPP / "val5_stop_frequency.png", "Applied vs Survey — Stop Frequency",
  "Stops per tour-half by purpose and direction (outbound/inbound); applied reproduces the survey's chaining rates."),
 (MAPP / "val6_stop_purpose.png", "Applied vs Survey — Stop Purpose (Activity Mix)",
  "Activity mix of intermediate stops by tour half; escort/drop-off skews outbound and shopping inbound, as in the survey."),
 (MAPP / "val7_stop_detour.png", "Applied vs Survey — Stop Detour",
  "Extra travel time added by an intermediate stop; applied median ≈19 min vs survey ≈14 (detours run slightly longer)."),
]
for path, ttl, ex in FIGS:
    big_figure(ttl, path, ex)

# ── summary table ──────────────────────────────────────────────────────────
s = slide(); title(s, "Tour-Based MITO — Summary")
rows = [
 ("Aspect", "Result"),
 ("Data / estimation", "RTS 2017–18 (163k trips, 42k persons); Apollo + Weibull AFT, survey-weighted, robust SE"),
 ("Estimation fit", "McFadden ρ² 0.27–0.81 across 9 components; all 27 runs converged"),
 ("Out-of-sample (held-out 20%)", "every component beats chance; DAP 0.79, mode 0.56–0.66, destination up to 19× chance"),
 ("Aggregate shares (OOS)", "predicted vs observed within 0.3–2.2 pp MAE"),
 ("Trip generation (applied)", "3.11 legs/person vs RTS 3.13; 1.13 tours/person vs 1.14"),
 ("Trip length (applied, mi)", "12.4 / 4.5 / 6.6 (HBW/HBS/HBO) = RTS exactly"),
 ("Policy readiness", "VOT-anchored cost coeff (~$22/hr) → I-695 pricing elasticity; cleared for MATSim"),
]
tbl = s.shapes.add_table(len(rows), 2, MARGIN, Inches(1.4), W - 2 * MARGIN, Inches(5.0)).table
tbl.columns[0].width = Inches(4.4); tbl.columns[1].width = Inches(7.7)
tblPr = tbl._tbl.tblPr
for child in list(tblPr):
    tblPr.remove(child)
for j, c in enumerate(rows[0]):
    cell = tbl.cell(0, j); cell.text = c
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    rr = cell.text_frame.paragraphs[0].runs[0]
    rr.font.bold = True; rr.font.size = Pt(13); rr.font.name = "Times New Roman"; rr.font.color.rgb = BLACK
for i in range(1, len(rows)):
    for j, c in enumerate(rows[i]):
        cell = tbl.cell(i, j); cell.text = c
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rr = cell.text_frame.paragraphs[0].runs[0]
        rr.font.size = Pt(11.5); rr.font.name = "Times New Roman"; rr.font.color.rgb = BLACK

prs.save(str(DECK))
print(f"appended MITO section → {DECK}  (now {len(prs.slides._sldIdLst)} slides)")
