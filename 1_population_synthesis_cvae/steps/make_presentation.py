#!/usr/bin/env python3
"""Build an academic PPTX summarizing the Updated VAE population-synthesis work.
Embeds the publication-quality figures from outputs/figures + outputs/09_study_areas.
Output → Updated_VAE_Presentation.pptx
"""
from __future__ import annotations
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
FT = ROOT / "outputs" / "figures" / "training"
FV = ROOT / "outputs" / "figures" / "validation"
SA = ROOT / "outputs" / "09_study_areas"
OUT = ROOT / "Updated_VAE_Presentation.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x5F); BLUE = RGBColor(0x2C, 0x6F, 0xBB)
GREY = RGBColor(0x44, 0x44, 0x44); RED = RGBColor(0xC4, 0x42, 0x3A); GREEN = RGBColor(0x2E, 0x7D, 0x32)
W, H = Inches(13.333), Inches(7.5)

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    return tf


def header(slide, title, sub=None):
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.12)); bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = _box(slide, Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.9))
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(14); rr.font.color.rgb = GREY; rr.font.italic = True


def bullets(slide, items, left=Inches(0.7), top=Inches(1.4), width=Inches(12), size=18):
    tf = _box(slide, left, top, width, Inches(5.6))
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 3 * lvl); r.font.color.rgb = GREY if lvl else NAVY
        p.space_after = Pt(6)


def image(slide, path, left, top, width=None, height=None):
    p = Path(path)
    if not p.exists():
        _box(slide, left, top, Inches(4), Inches(0.4)).paragraphs[0].add_run().text = f"[missing {p.name}]"
        return
    kw = {}
    if width: kw["width"] = width
    if height: kw["height"] = height
    slide.shapes.add_picture(str(p), left, top, **kw)


def caption(slide, txt, left, top, width):
    tf = _box(slide, left, top, width, Inches(0.5))
    r = tf.paragraphs[0].add_run(); r.text = txt; r.font.size = Pt(11); r.font.italic = True
    r.font.color.rgb = GREY; tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def newslide():
    return prs.slides.add_slide(BLANK)


# ── 1. Title ──────────────────────────────────────────────────────────────
s = newslide()
bg = s.shapes.add_shape(1, 0, 0, W, H); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tf = _box(s, Inches(1), Inches(2.1), Inches(11.3), Inches(2.5))
r = tf.paragraphs[0].add_run()
r.text = "A Simple Conditional VAE for Synthetic Population Synthesis"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)
p = tf.add_paragraph(); r = p.add_run()
r.text = "ACS PUMS 2016 → SILO / MITO, Baltimore–Washington MSTM region"
r.font.size = Pt(22); r.font.color.rgb = RGBColor(0xBF, 0xD3, 0xE8)
p = tf.add_paragraph(); p.space_before = Pt(20); r = p.add_run()
r.text = "Methodology · Validation framework · Honest limitations"
r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = RGBColor(0x9F, 0xB8, 0xD0)

# ── 2. Background & motivation ────────────────────────────────────────────
s = newslide(); header(s, "Background & Motivation")
bullets(s, [
 ("Transport microsimulation pipeline: VAE → SILO (land use) → MITO (travel demand) → MATSim (assignment).", 0),
 ("SILO/MITO need a complete synthetic base-year population: households + persons + dwellings + jobs.", 0),
 ("Source: U.S. Census ACS PUMS 2016 5-year microdata for the 96-PUMA MSTM region (DE/DC/MD/PA/VA/WV).", 0),
 ("Prior model (v6): a hierarchical CVAE with a multi-head income model.", 0),
 ("Two diagnosed failures motivating a rebuild:", 0),
 ("Heavy-tailed household income underfit — P95 ≈ −11% (the rich tail too thin).", 1),
 ("Posterior collapse — only 3–4 of 96 latent dimensions active.", 1),
 ("Goal: a SIMPLE, transparent CVAE that fixes income and the collapse without patches or over-engineering.", 0),
])

# ── 3. Data ───────────────────────────────────────────────────────────────
s = newslide(); header(s, "Data — ACS PUMS 2016 (5-year)")
bullets(s, [
 ("253,029 households / 626,201 persons (weighted to 4.70M HH / 12.49M persons).", 0),
 ("Region: 96 PUMAs across 6 states; base year 2016 (incomes via ACS ADJINC).", 0),
 ("Modeled household variables: dwellingType, tenure, autos, income.", 0),
 ("Modeled person variables: age, gender, race, occupation, driver's-license, relationship, nationality, income.", 0),
 ("All recodes verified against the official 2012–2016 Census PUMS Data Dictionary + actual data ranges.", 0),
 ("Empirical income facts driving the design:", 0),
 ("Continuous to $2.2M, no native top-code; sub-Pareto / thinning upper tail (Hill α rises 2.7→7).", 1),
 ("~28% structural-zero person incomes (non-earners); strong spatial variation.", 1),
])

# ── 4. The simple-CVAE design ─────────────────────────────────────────────
s = newslide(); header(s, "Approach — A Simple Conditional VAE", "design principles: simplicity, no patches")
bullets(s, [
 ("Plain CVAE (Borysov 2019 style): flat one-hot record → MLP encoder → Gaussian latent → MLP decoder → softmax heads.", 0),
 ("PUMA is the only embedding; every other categorical is one-hot.", 0),
 ("No autoregressive decoder, no Set Transformer, no in-network income head.", 0),
 ("Income as a categorical bin + empirical within-(PUMA, bin) draw at generation (open top bin keeps $1M+).", 0),
 ("Person income via earner bin (28% zeros) + exact reconciliation Σperson = household income.", 0),
 ("Latent dim 24 with KL free-bits → prevents the collapse that crippled v6.", 0),
 ("Per-PUMA marginal-JSD loss (w=2.0) to tighten spatial marginals; constrained decoding → 0 structural zeros.", 0),
])

# ── 5. Income representation (the key idea) ───────────────────────────────
s = newslide(); header(s, "Key Idea — Binned Income + Within-Bin Empirical Draw")
bullets(s, [
 ("Heavy/thinning tail cannot be captured by a Gaussian decoder (provable for light-tailed bases).", 0),
 ("Instead: the VAE predicts an income BIN; dollars are drawn from real records in that (PUMA, bin) cell.", 0),
 ("The open top bin returns genuine $1M–$2.2M incomes — no parametric tail head, no clip.", 0),
], width=Inches(6))
image(s, FV / "F5_income_distribution.png", Inches(6.7), Inches(1.5), width=Inches(6.2))
caption(s, "Household income: VAE generated vs ACS test ($0–400k).", Inches(6.7), Inches(6.6), Inches(6.2))

# ── 6. Validation framework ───────────────────────────────────────────────
s = newslide(); header(s, "Validation Framework", "held-out, floor-grounded, journal-style")
bullets(s, [
 ("Per-PUMA, household-level 80/10/10 train/validation/test split of the 2016 PUMS (no out-of-vintage data).", 0),
 ("Train fits; validation = early stopping; TEST = held-out honest evaluation.", 0),
 ("12-category suite: totals, marginals (F1), joints by order (F8), association (S3), income (F5),", 0),
 ("household structure (S6), spatial (per-PUMA), structural zeros, sampling zeros, memorization, coherence.", 1),
 ("Identifiability floors: half-split bootstrap of the test set bounds the best any model could achieve.", 0),
 ("Metrics: Total-Variation & SRMSE with bootstrap CIs; effect sizes, not p-values (Müller & Axhausen).", 0),
])

# ── 7. Result: income tail (the win) ──────────────────────────────────────
s = newslide(); header(s, "Result — Income Tail Recovered (the main win)")
image(s, FV / "F5_income_tail_ccdf.png", Inches(0.5), Inches(1.4), width=Inches(7.2))
caption(s, "Log-log survival: generated tracks ACS test across the full range, incl. $1M+ tail.", Inches(0.5), Inches(6.5), Inches(7.2))
bullets(s, [
 ("P95 bias −2.6% (v6: −11%).", 0),
 ("P99 bias −4.2%; max $2.2M.", 0),
 (">$300k share 3.6% vs 3.8% truth.", 0),
 ("Income marginal & per-PUMA spatial both at/below the identifiability floor — effectively solved.", 0),
], left=Inches(8.0), top=Inches(1.7), width=Inches(5.0), size=16)

# ── 8. Result: marginals + nationality ────────────────────────────────────
s = newslide(); header(s, "Result — Marginal Fidelity + Nationality")
image(s, FV / "F1_marginal_TV_summary.png", Inches(0.4), Inches(1.5), width=Inches(7.6))
caption(s, "Per-attribute Total-Variation vs held-out test (red = 0.03 ref).", Inches(0.4), Inches(6.4), Inches(7.6))
bullets(s, [
 ("At/below floor: income, gender, dwelling, autos, nationality.", 0),
 ("Nationality added to the VAE: captures the spatial gradient —", 0),
 ("NoVA 24% foreign-born vs Baltimore 9% (truth 27% / 10%).", 1),
 ("Above floor (bounded): occupation, age, relationship, race.", 0),
], left=Inches(8.2), top=Inches(1.7), width=Inches(4.9), size=15)

# ── 9. Result: joints ─────────────────────────────────────────────────────
s = newslide(); header(s, "Result — Joint Distributions")
image(s, FV / "F8_joint_srmse_by_order.png", Inches(0.6), Inches(1.6), width=Inches(5.6))
image(s, FV / "F3_joint_age_bin_x_income_bin.png", Inches(6.4), Inches(2.0), width=Inches(6.5))
caption(s, "SRMSE by interaction order.", Inches(0.6), Inches(6.5), Inches(5.6))
caption(s, "Age × income joint: truth / generated / difference.", Inches(6.4), Inches(5.8), Inches(6.5))

# ── 10. Training diagnostics: no collapse ─────────────────────────────────
s = newslide(); header(s, "Training — No Posterior Collapse, No Over/Underfit")
image(s, FT / "T3_kl_active_dims.png", Inches(0.5), Inches(1.5), width=Inches(6.2))
image(s, FT / "T1_elbo_curves.png", Inches(7.0), Inches(1.5), width=Inches(6.0))
caption(s, "All 24 latent dims active (free-bits) — vs v6's 3–4/96.", Inches(0.5), Inches(6.5), Inches(6.2))
caption(s, "Train/val ELBO track closely → good fit.", Inches(7.0), Inches(6.5), Inches(6.0))

# ── 11. Study areas ───────────────────────────────────────────────────────
s = newslide(); header(s, "Study-Area Validation — Baltimore & Northern Virginia")
image(s, SA / "nova_county_TV_heatmap.png", Inches(0.4), Inches(1.5), width=Inches(7.4))
bullets(s, [
 ("County = its PUMA set (zone-system majority).", 0),
 ("Income medians within a few %:", 0),
 ("Baltimore −1.3%, NoVA −1.5% overall.", 1),
 ("Per-county marginal TV 0.023–0.042.", 0),
 ("Nationality gradient reproduced across", 0),
 ("Fairfax / Loudoun / Arlington / PW.", 1),
], left=Inches(8.1), top=Inches(1.7), width=Inches(5.0), size=15)
caption(s, "NoVA: per-county marginal TV heatmap.", Inches(0.4), Inches(6.5), Inches(7.4))

# ── 12. Gap analysis vs floor (table) ─────────────────────────────────────
s = newslide(); header(s, "Gap Analysis — Real Gaps vs Identifiability Floor")
rows = [
 ("Attribute", "gen TV", "floor", "status"),
 ("income (HH) / per-PUMA", "0.033 / 0.30", "0.029 / 0.53", "AT/BELOW FLOOR — solved"),
 ("nationality / gender / dwelling / autos", "≤0.014", "≈", "AT FLOOR — solved"),
 ("occupation", "0.040", "0.008", "above floor (bounded)"),
 ("age", "0.038", "0.011", "above floor (bounded)"),
 ("relationship", "0.033", "0.007", "above floor (bounded)"),
 ("race", "0.026", "0.008", "above floor (bounded)"),
 ("couple age gap", "7.55 y", "3.69 y", "architectural limit"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.6)).table
tbl.columns[0].width = Inches(4.6); tbl.columns[1].width = Inches(2.2)
tbl.columns[2].width = Inches(2.0); tbl.columns[3].width = Inches(3.3)
for j, c in enumerate(rows[0]):
    cell = tbl.cell(0, j); cell.text = c
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
    cell.text_frame.paragraphs[0].runs[0].font.bold = True; cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
for i in range(1, len(rows)):
    solved = "FLOOR" in rows[i][3]
    for j, c in enumerate(rows[i]):
        cell = tbl.cell(i, j); cell.text = c
        run = cell.text_frame.paragraphs[0].runs[0]; run.font.size = Pt(12)
        run.font.color.rgb = GREEN if (j == 3 and solved) else (RED if j == 3 else GREY)

# ── 13. Limitations & diagnosis ───────────────────────────────────────────
s = newslide(); header(s, "Limitations — Diagnosed & Bounded")
image(s, FV / "S6_couple_age_gap.png", Inches(7.0), Inches(1.5), width=Inches(6.0))
bullets(s, [
 ("Residual gaps (occupation/age/relationship/race, couple-gap) share one cause:", 0),
 ("weak within-household coupling — person heads are conditionally independent given z.", 1),
 ("Couple age-gap is a model BIN-level error (+3.3y), not the within-bin draw (+0.6y).", 0),
 ("Ablations tried (all within constraints):", 0),
 ("larger latent → worse; lower β → worse; deeper decoder → wash; marginal-JSD → modest, kept.", 1),
 ("Only effective fixes are out of scope by design choice:", 0),
 ("autoregressive/cross-person decoder, or post-hoc IPU calibration (→ age 0.044→0.009).", 1),
], left=Inches(0.6), top=Inches(1.5), width=Inches(6.2), size=15)
caption(s, "Couple age gap: gen 7.6y vs test 3.7y (the bounded limit).", Inches(7.0), Inches(6.5), Inches(6.0))

# ── 14. Conclusion ────────────────────────────────────────────────────────
s = newslide(); header(s, "Conclusion & Deliverable")
bullets(s, [
 ("A simple, transparent CVAE that FIXES the two v6 failures:", 0),
 ("income tail recovered (P95 −2.6% vs −11%) via binning + within-bin empirical draw;", 1),
 ("no posterior collapse (24/24 dims) via free-bits.", 1),
 ("Held-out validation: marginals largely at floor; structural zeros 0; coherence 100%; nationality gradient captured.", 0),
 ("Honest, floor-grounded accounting: income/nationality/spatial solved; person-coupling marginals are the bounded ceiling.", 0),
 ("Deliverable: 4.70M-household SILO population (hh/pp/dd/jj_2016.csv, verified schema) + 35 publication figures.", 0),
 ("Fully reproducible: numbered pipeline steps 00→09, self-contained inputs.", 0),
], size=17)

prs.save(str(OUT))
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
print(f"file size: {OUT.stat().st_size/1e6:.1f} MB")
