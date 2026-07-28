#!/usr/bin/env python3
"""APPEND a SILO section to the (user-edited) VAE deck WITHOUT regenerating it.
Opens VAE_Architecture_and_Validation.pptx, appends SILO calibration/validation
methodology (plain diagrams) + Maryland-only year-to-year validation (one clean
figure per slide, VAE style), and saves in place. Same plain look as the VAE deck.

IMPORTANT: close PowerPoint first (the deck must not be open), or your save will
clobber these slides.
"""
from __future__ import annotations
import sys
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
# optional override: python append_silo_slides.py <in_out.pptx>  (used for a dry-run on a copy)
DECK = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "VAE_Architecture_and_Validation.pptx"
SILO = ROOT.parent / "Updated SILO"
VFIG = SILO / "validation" / "by_year_acs_calib5" / "figures"
MD23 = SILO / "validation" / "by_year_acs_calib5" / "2023" / "MD"

BLACK = RGBColor(0, 0, 0); GREY = RGBColor(0x55, 0x55, 0x55); LINE = RGBColor(0x33, 0x33, 0x33)

prs = Presentation(str(DECK))
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
MARGIN = Inches(0.6)


def slide():
    return prs.slides.add_slide(BLANK)


def _tb(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
    return tf


def title(s, text, sub=None, size=24):
    tf = _tb(s, MARGIN, Inches(0.22), W - 2 * MARGIN, Inches(1.0))
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    rule_y = Inches(0.78)
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(13); rr.font.italic = True; rr.font.name = "Times New Roman"; rr.font.color.rgb = GREY
        rule_y = Inches(1.12)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, rule_y, W - 2 * MARGIN, Pt(1.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    return rule_y


def bullets(s, items, left=MARGIN, top=Inches(1.5), width=None, size=15):
    width = width or (W - 2 * MARGIN)
    tf = _tb(s, left, top, width, H - top - Inches(0.4))
    for i, it in enumerate(items):
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(5)
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size - 2 * lvl); r.font.name = "Times New Roman"
        r.font.color.rgb = BLACK if lvl == 0 else GREY


def block(s, l, t, w, h, text, fontsize=12, bold=False):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(255, 255, 255)
    b.line.color.rgb = BLACK; b.line.width = Pt(1); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(fontsize); r.font.bold = bold
    r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    return b


def arrow(s, l, t, w, kind=MSO_SHAPE.RIGHT_ARROW, h=Inches(0.22)):
    a = s.shapes.add_shape(kind, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88)
    a.line.fill.background(); a.shadow.inherit = False


def big_figure(title_txt, path, explain):
    s = slide(); rule_y = title(s, title_txt, explain, size=22)
    p = Path(path)
    if not p.exists():
        _tb(s, MARGIN, Inches(3), Inches(6), Inches(0.4)).paragraphs[0].add_run().text = f"[missing {p.name}]"
        return s
    iw, ih = Image.open(str(p)).size; ar = iw / ih
    box_l = Inches(0.25); box_t = int(rule_y) + Inches(0.12)
    box_w = W - Inches(0.5); box_h = H - box_t - Inches(0.15)
    if ar > box_w / box_h:
        w, h = int(box_w), int(box_w / ar)
    else:
        w, h = int(box_h * ar), int(box_h)
    l = int(box_l + (box_w - w) / 2); t = int(box_t + (box_h - h) / 2)
    s.shapes.add_picture(str(p), l, t, width=w, height=h)
    return s


def section(text, subtitle=None):
    s = slide()
    tf = _tb(s, MARGIN, Inches(2.8), W - 2 * MARGIN, Inches(1.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.name = "Times New Roman"; r.font.color.rgb = BLACK
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.3), Inches(4.0), Inches(4.7), Pt(1.4))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    if subtitle:
        tf2 = _tb(s, MARGIN, Inches(4.2), W - 2 * MARGIN, Inches(1.2))
        p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = subtitle
        r.font.size = Pt(14); r.font.italic = True; r.font.name = "Times New Roman"; r.font.color.rgb = GREY
    return s


# ══════════════════════════════════════════════════════════════════════════
# SILO — section divider
# ══════════════════════════════════════════════════════════════════════════
section("SILO — Land-Use Microsimulation",
        "Calibration & validation methodology, and Maryland year-to-year forecast validation vs ACS PUMS")

# ── 1. Overview / pipeline ─────────────────────────────────────────────────
s = slide(); title(s, "SILO — Overview",
                   "a closed-loop land-use microsimulation: ages the VAE 2016 base population forward, one year at a time, to 2023")
y = Inches(1.75); bh = Inches(1.05)
block(s, Inches(0.55), y, Inches(2.6), bh, "VAE base-year population 2016\n(households · persons · dwellings · jobs)", 11)
arrow(s, Inches(3.22), Inches(2.15), Inches(0.5))
block(s, Inches(3.78), y, Inches(2.9), bh, "SILO annual microsimulation\n2016 → 2023\n(demographic · employment · real-estate · mobility)", 11, True)
arrow(s, Inches(6.74), Inches(2.15), Inches(0.5))
block(s, Inches(7.3), y, Inches(2.4), bh, "Yearly synthetic population\n(hh · pp · dd · jj per year)", 11)
arrow(s, Inches(9.76), Inches(2.15), Inches(0.5))
block(s, Inches(10.32), y, Inches(2.45), bh, "Validate vs ACS PUMS\n5-year (per state, per year)", 11)
bullets(s, [
 ("Region: 6 states (DE, DC, MD, PA, VA, WV) across 96 MSTM PUMAs; this deck focuses the validation on Maryland.", 0),
 ("Closed model: no exogenous re-seeding — every household and person is aged forward by micro-event and "
  "annual-update models, so forecast accuracy is an honest test of the calibrated dynamics.", 0),
 ("Base year (2016) comes straight from the VAE synthetic population; SILO is never re-fit to it.", 0),
 ("Each simulated year emits a full population that is scored against ACS PUMS 5-year.", 0),
], top=Inches(3.15), size=15)

# ── 2. Annual simulation loop (models) ─────────────────────────────────────
s = slide(); title(s, "SILO — Annual Simulation Loop (per year, 2016 → 2023)",
                   "every household & person passes through micro-event models, then annual market-update models")
lx, rx = Inches(0.7), Inches(7.0); cw = Inches(5.6); by = Inches(1.9); bh = Inches(0.62); gap = Inches(0.72)
# column headers
for x, htxt in ((lx, "EVENT models  (micro-events)"), (rx, "ANNUAL update models")):
    tf = _tb(s, x, Inches(1.45), cw, Inches(0.4)); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = htxt; r.font.size = Pt(13); r.font.italic = True
    r.font.name = "Times New Roman"; r.font.color.rgb = GREY
EVENTS = ["birth · birthday · death", "marriage · divorce", "leave-parent · driver's license",
          "education · employment", "construction · demolition", "renovation · move · migration"]
ANNUAL = ["job-market update", "real-estate pricing", "construction overwrite",
          "auto-ownership (calibrated, wired in)", "income adjustment"]
for i, txt in enumerate(EVENTS):
    block(s, lx, int(by) + i * int(gap), cw, bh, txt, 12)
for i, txt in enumerate(ANNUAL):
    block(s, rx, int(by) + i * int(gap), cw, bh, txt, 12)
bullets(s, [
 ("Event models fire stochastic life-events per person/household; annual models re-clear the housing "
  "and job markets and re-price/adjust each year.", 0),
], top=Inches(6.5), size=13)

# ── 3. Calibration & validation methodology (the loop) ─────────────────────
s = slide(); title(s, "SILO — Calibration & Validation Methodology",
                   "calibrate in-sample (2016–2020), FREEZE, then forecast out-of-sample (2021–2023)")
y1 = Inches(1.7); bh = Inches(0.95)
block(s, Inches(0.55), y1, Inches(2.7), bh, "Run SILO\n2016 – 2020", 12, True)
arrow(s, Inches(3.32), Inches(2.05), Inches(0.4))
block(s, Inches(3.8), y1, Inches(3.0), bh, "Validate vs ACS\n(in-sample fit)", 12)
arrow(s, Inches(6.9), Inches(2.05), Inches(0.4))
block(s, Inches(7.38), y1, Inches(5.4), bh, "Adjust per-state levers\n(birth · marriage · income · auto)", 12)
# iterate arrow (down-left loop label)
tf = _tb(s, Inches(3.8), Inches(2.75), Inches(5.0), Inches(0.35))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "↑  iterate until the 2016–2020 fit converges"
r.font.size = Pt(12); r.font.italic = True; r.font.name = "Times New Roman"; r.font.color.rgb = GREY
y2 = Inches(3.35)
block(s, Inches(0.55), y2, Inches(2.7), bh, "FREEZE levers", 12, True)
arrow(s, Inches(3.32), Inches(3.7), Inches(0.4))
block(s, Inches(3.8), y2, Inches(3.0), bh, "Forecast SILO\n2021 – 2023", 12, True)
arrow(s, Inches(6.9), Inches(3.7), Inches(0.4))
block(s, Inches(7.38), y2, Inches(5.4), bh, "Validate vs ACS\n(OUT-OF-SAMPLE forecast skill)", 12)
bullets(s, [
 ("Acceptance is PER-BIN: every category of every variable within ±5 pp of ACS (stricter than an aggregate "
  "Total-Variation score).", 0),
 ("Per-state levers: birth, marriage, income and an auto-ownership ASC that self-calibrates to 2016 shares, then freezes.", 0),
 ("Composition re-anchoring (race · household-size · income) corrects closed-model drift — migration is "
  "composition-neutral, so shares would otherwise freeze while ACS moves.", 0),
 ("2021–2023 is never used for calibration → the out-of-sample fit is the honest forecast-skill number.", 0),
], top=Inches(4.5), size=14)

# ── 4. Maryland year-to-year (one clean figure) ────────────────────────────
big_figure("Maryland — Year-to-Year Validation (2016 → 2023)",
           VFIG / "md_year_to_year.png",
           "Total Variation of SILO vs ACS PUMS for each variable; shaded = out-of-sample forecast. "
           "All variables stay low — age is the largest (~0.06); occupation/autos/race next.")

# ── 5. Maryland per-variable validation, 2023 out-of-sample (VAE style) ─────
MDVARS = [
 ("pp_age_bin.png",      "Maryland Validation 2023 — Age"),
 ("pp_occ_silo.png",     "Maryland Validation 2023 — Occupation"),
 ("hh_autos.png",        "Maryland Validation 2023 — Autos per Household"),
 ("hh_hhSize.png",       "Maryland Validation 2023 — Household Size"),
 ("hh_dwellingType.png", "Maryland Validation 2023 — Dwelling Type"),
 ("pp_race4.png",        "Maryland Validation 2023 — Race"),
 ("hh_hh_inc9.png",      "Maryland Validation 2023 — Household Income"),
 ("pp_gender.png",       "Maryland Validation 2023 — Gender"),
]
EXPL = ("Top: ACS (observed) vs SILO (model) shares. Bottom: SILO − ACS per-bin gap; "
        "green band = ±5 pp acceptance. 2023 is fully out-of-sample.")
for fn, ttl in MDVARS:
    big_figure(ttl, MD23 / fn, EXPL)

prs.save(str(DECK))
print(f"appended SILO section → {DECK}  (now {len(prs.slides._sldIdLst)} slides)")
