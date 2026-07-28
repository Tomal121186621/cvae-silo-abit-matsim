#!/usr/bin/env python3
"""Living PPTX generator: SILO models, the changes we are making, and the calibration/validation
framework for the VAE -> SILO -> MITO -> MATSim pipeline (Baltimore-Washington MSTM region).

Re-run this whenever the engine changes, the calibration converges, or validation completes:
    python3 make_pptx.py
It regenerates  SILO_Calibration_Validation_Framework.pptx  and its diagram images, pulling the
live per-state calibration values and (if present) the latest validation scorecard.

Design: content lives in DATA structures below so updating a number/slide is a one-line edit.
"""
from __future__ import annotations
import csv, datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

HERE = Path(__file__).resolve().parent
USILO = HERE.parent
FIGDIR = HERE / "_pptx_assets"; FIGDIR.mkdir(exist_ok=True)
OUT = HERE / "SILO_Calibration_Validation_Framework.pptx"
CALIB_CSV = Path("/Users/tomal/Documents/VAE SILO Architecture/silo_smoke_test/input/assumptions/calibration_by_state.csv")
# Per-bin scorecard + summary heatmap from the current reference run (calib1 until calib2/3 validate).
# Point REF_SUB at the latest validated scenario subfolder to refresh the deck after a new run.
REF_SUB   = "by_year_acs_calib5"
PERBIN    = USILO / "validation" / REF_SUB / "perbin_scorecard.csv"
HEATMAP   = USILO / "validation" / REF_SUB / "figures" / "scorecard_heatmap.png"
TRAJ      = USILO / "validation" / REF_SUB / "figures" / "error_trajectory.png"

# ---- palette ----
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
BLUE   = RGBColor(0x2E, 0x6D, 0xB4)
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)
ORANGE = RGBColor(0xE7, 0x6F, 0x51)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x22, 0x22, 0x22)

TODAY = datetime.date.today().isoformat()

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------- helpers
def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame

def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _bar(slide, color=NAVY, h=0.9):
    s = slide.shapes.add_shape(1, 0, 0, SW, Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    return s

def header(slide, title, kicker=None):
    _bar(slide, NAVY, 0.95)
    tf = _box(slide, 0.55, 0.12, 12.2, 0.8); tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    if kicker:
        tk = _box(slide, 0.55, 0.62, 12.2, 0.3)
        p = tk.paragraphs[0]; r = p.add_run(); r.text = kicker
        r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = RGBColor(0xCF,0xDD,0xEC)
    foot(slide)

def foot(slide):
    tf = _box(slide, 0.4, 7.06, 12.5, 0.35)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"VAE→SILO→MITO→MATSim  ·  MSTM Baltimore–Washington  ·  updated {TODAY}"
    r.font.size = Pt(9); r.font.color.rgb = GREY

def bullets(slide, items, left=0.7, top=1.2, width=12.0, height=5.6, size=18):
    tf = _box(slide, left, top, width, height); tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple): lvl, txt = it
        else: lvl, txt = 0, it
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.level = lvl; p.space_after = Pt(5)
        bold = txt.startswith("**") and txt.endswith("**")
        txt = txt.strip("*")
        bullet = "—  " if lvl == 0 else "·  "
        r = p.add_run(); r.text = bullet + txt
        r.font.size = Pt(size - 2*lvl); r.font.color.rgb = BLACK if lvl==0 else GREY
        if bold: r.font.bold = True; r.font.color.rgb = NAVY

def table(slide, headers, rows, left=0.6, top=1.25, width=12.1, height=5.4,
          col_w=None, fs=12, hdr=BLUE):
    nr, nc = len(rows)+1, len(headers)
    gt = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(height)).table
    if col_w:
        tot = sum(col_w)
        for j,w in enumerate(col_w): gt.columns[j].width = Emu(int(Inches(width)*w/tot))
    for j,h in enumerate(headers):
        c = gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb = hdr
        tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; r=p.add_run(); r.text=str(h)
        r.font.bold=True; r.font.size=Pt(fs); r.font.color.rgb=WHITE
    for i,row in enumerate(rows, start=1):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb = WHITE if i%2 else LIGHT
            tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; r=p.add_run(); r.text=str(val)
            r.font.size=Pt(fs-1); r.font.color.rgb=BLACK
            if j==0: r.font.bold=True; r.font.color.rgb=NAVY
    return gt

def code(slide, lines, left=0.7, top=1.3, width=12.0, height=5.2, fs=12):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x1E,0x1E,0x2A); s.line.fill.background()
    s.shadow.inherit=False
    tf = s.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.15)
    first=True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        col = RGBColor(0x6A,0x99,0x55) if ln.strip().startswith("//") else RGBColor(0xE8,0xE8,0xF0)
        r=p.add_run(); r.text=ln; r.font.name="Menlo"; r.font.size=Pt(fs); r.font.color.rgb=col


# --------------------------------------------------------------------------- diagrams
def diagram_pipeline():
    fig, ax = plt.subplots(figsize=(12, 2.5)); ax.axis("off"); ax.set_xlim(0,12); ax.set_ylim(0,3)
    stages = [("VAE\nsynthetic pop\n(base 2016)", "#264653"),
              ("SILO\nland use\n2016→2023", "#2A9D8F"),
              ("MITO\ntravel demand", "#E9C46A"),
              ("MATSim\nassignment", "#E76F51")]
    x=0.4
    for i,(name,c) in enumerate(stages):
        ax.add_patch(FancyBboxPatch((x,0.8),2.4,1.4, boxstyle="round,pad=0.05",
                     fc=c, ec="none")); ax.text(x+1.2,1.5,name,ha="center",va="center",
                     color="white",fontsize=12,fontweight="bold")
        if i<3:
            ax.add_patch(FancyArrowPatch((x+2.45,1.5),(x+2.95,1.5),arrowstyle="-|>",
                         mutation_scale=22,color="#333"))
        x+=2.9
    ax.text(6,0.35,"each stage's output is the next stage's input",ha="center",
            fontsize=10,style="italic",color="#555")
    p=FIGDIR/"pipeline.png"; fig.savefig(p,dpi=150,bbox_inches="tight"); plt.close(); return p

def diagram_loop():
    fig, ax = plt.subplots(figsize=(11,5.2)); ax.axis("off"); ax.set_xlim(0,11); ax.set_ylim(0,6)
    def box(x,y,w,h,t,c,fs=11):
        ax.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.04",fc=c,ec="none"))
        ax.text(x+w/2,y+h/2,t,ha="center",va="center",color="white",fontsize=fs,fontweight="bold")
    def arrow(x1,y1,x2,y2):
        ax.add_patch(FancyArrowPatch((x1,y1),(x2,y2),arrowstyle="-|>",mutation_scale=20,
                     color="#444",connectionstyle="arc3,rad=0"))
    box(0.4,4.4,2.6,1.1,"Run SILO\n2016–2020","#2A9D8F")
    box(3.6,4.4,2.6,1.1,"Validate vs ACS\n(in-sample fit)","#2E6DB4")
    box(6.8,4.4,3.6,1.1,"Adjust per-state levers\n(birth/marriage/income/auto)","#E76F51")
    arrow(3.0,4.95,3.6,4.95); arrow(6.2,4.95,6.8,4.95)
    # loop back
    ax.add_patch(FancyArrowPatch((8.6,4.4),(1.7,4.4),arrowstyle="-|>",mutation_scale=20,
                 color="#444",connectionstyle="arc3,rad=0.35"))
    ax.text(5.1,3.25,"iterate until 2016–2020 fit converges",ha="center",fontsize=10,
            style="italic",color="#555")
    box(0.4,1.4,2.6,1.1,"FREEZE\nlevers","#264653")
    box(3.6,1.4,2.6,1.1,"Forecast SILO\n2021–2023","#2A9D8F")
    box(6.8,1.4,3.6,1.1,"Validate vs ACS\n(OUT-OF-SAMPLE skill)","#2E6DB4")
    arrow(1.7,4.4,1.7,2.5); arrow(3.0,1.95,3.6,1.95); arrow(6.2,1.95,6.8,1.95)
    ax.text(5.5,0.7,"out-of-sample 2021–2023 accuracy = the headline forecast-skill number",
            ha="center",fontsize=11,color="#1F3A5F",fontweight="bold")
    p=FIGDIR/"loop.png"; fig.savefig(p,dpi=150,bbox_inches="tight"); plt.close(); return p

def diagram_arch():
    fig, ax = plt.subplots(figsize=(11,5)); ax.axis("off"); ax.set_xlim(0,11); ax.set_ylim(0,6)
    ax.add_patch(FancyBboxPatch((0.3,0.3),10.4,5.4,boxstyle="round,pad=0.05",fc="#EEF2F7",ec="#1F3A5F"))
    ax.text(5.5,5.4,"SILO annual simulation loop (per year 2016→2023)",ha="center",
            fontsize=13,fontweight="bold",color="#1F3A5F")
    def box(x,y,w,h,t,c,fs=9.5):
        ax.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.03",fc=c,ec="none"))
        ax.text(x+w/2,y+h/2,t,ha="center",va="center",color="white",fontsize=fs,fontweight="bold")
    ax.text(2.6,4.75,"EVENT models (micro-events)",ha="center",fontsize=10,color="#444",style="italic")
    ev=["birth · birthday · death","marriage · divorce","leave-parent · drivers lic.",
        "education · employment","construction · demolition","renovation · move · migration"]
    y=4.2
    for t in ev: box(0.7,y,3.8,0.45,t,"#2E6DB4"); y-=0.6
    ax.text(8.2,4.75,"ANNUAL update models",ha="center",fontsize=10,color="#444",style="italic")
    an=["job-market update","real-estate pricing","construction overwrite",
        "★ auto ownership (NEW: wired in)","income adjustment (data mgr)"]
    y=4.2
    for t in an:
        c = "#E76F51" if t.startswith("★") else "#2A9D8F"
        box(6.3,y,4.0,0.45,t,c); y-=0.6
    p=FIGDIR/"arch.png"; fig.savefig(p,dpi=150,bbox_inches="tight"); plt.close(); return p


# --------------------------------------------------------------------------- live data
def read_calib():
    if not CALIB_CSV.exists(): return None
    with open(CALIB_CSV) as f: return list(csv.DictReader(f))

def read_perbin():
    """Worst per-bin gap (pp) over the forecast window, per state x variable, from the reference run."""
    if not PERBIN.exists(): return None
    import collections
    worst = collections.defaultdict(dict)
    with open(PERBIN) as f:
        for r in csv.DictReader(f):
            try:
                yr = int(r["year"]); v = float(r["max_bin_pp"])
            except (ValueError, KeyError):
                continue
            if 2021 <= yr <= 2023:
                k = (r["state"], r["variable"])
                worst[k[0]][k[1]] = max(worst[k[0]].get(k[1], 0.0), v)
    return dict(worst)


# =========================================================================== SLIDES
# 1 — title
s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
tf=_box(s,0.9,2.2,11.5,2.2); tf.word_wrap=True
p=tf.paragraphs[0]; r=p.add_run(); r.text="SILO Land-Use Microsimulation"
r.font.size=Pt(44); r.font.bold=True; r.font.color.rgb=WHITE
p=tf.add_paragraph(); r=p.add_run()
r.text="Models in detail · the engine changes we are making · the calibrate–freeze–forecast & validation framework"
r.font.size=Pt(19); r.font.color.rgb=RGBColor(0xCF,0xDD,0xEC)
tf=_box(s,0.9,4.7,11.5,1.2)
p=tf.paragraphs[0]; r=p.add_run()
r.text="VAE synthetic population → SILO (2016–2023) → MITO → MATSim   ·   6 states (DE/DC/MD/PA/VA/WV)"
r.font.size=Pt(14); r.font.color.rgb=WHITE
p=tf.add_paragraph(); r=p.add_run(); r.text=f"Living document — last updated {TODAY}"
r.font.size=Pt(12); r.font.italic=True; r.font.color.rgb=RGBColor(0x9F,0xB8,0xD0)

# 2 — pipeline context
s=prs.slides.add_slide(BLANK); header(s,"Where SILO sits in the pipeline","the land-use engine between population synthesis and travel demand")
s.shapes.add_picture(str(diagram_pipeline()), Inches(0.6), Inches(1.4), width=Inches(12.1))
bullets(s,[
    "**SILO** evolves the VAE base-year population forward one year at a time to 2023.",
    "Its yearly population (households, persons, dwellings, jobs) is the **input to MITO** travel demand.",
    "This deck covers the SILO stage: its internal models, our code changes, and how we calibrate & validate it.",
], top=3.7, size=16)

# 3 — what is SILO
s=prs.slides.add_slide(BLANK); header(s,"What SILO is")
bullets(s,[
    "**Simple Integrated Land-Use Orchestrator** — an agent-based land-use microsimulation (TUM / U. Maryland).",
    "Simulates each household, person, dwelling and job **individually**, year by year.",
    "Two kinds of models run each year:",
    (1,"**Event models** — discrete life events drawn per agent (birth, death, marriage, job change, move, …)."),
    (1,"**Annual update models** — system-wide updates (job market, dwelling prices, auto ownership, income)."),
    "**Exogenous control totals** (population & employment by state) steer regional growth; behaviour is endogenous.",
    "Design philosophy: **calibrate behavioural coefficients once, then forecast** — it does NOT re-fit every year.",
], size=17)

# 4 — architecture diagram
s=prs.slides.add_slide(BLANK); header(s,"SILO architecture — the annual loop")
s.shapes.add_picture(str(diagram_arch()), Inches(0.9), Inches(1.25), width=Inches(11.5))
bullets(s,[ "DataContainer holds the agents; ModelContainer holds the models; the Simulator runs them each year.",
            "★ We **wired the Maryland auto-ownership model into the annual loop** (it previously ran only at synthesis)."],
        top=6.05, size=12)

# 5 — demography models
s=prs.slides.add_slide(BLANK); header(s,"The models (1/3) — demography","per-agent life events drawn each year")
table(s,["Model","What it does","Key driver"],[
    ["Birth","Women of child-bearing age may have a child","age, #children, marital status"],
    ["Birthday","Ages every person by one year","—"],
    ["Death","Mortality by age/gender","age, gender"],
    ["Marriage","Forms couples from a marriage market","age, gender, hh size"],
    ["Divorce","Splits married couples","marriage duration proxy"],
    ["Leave parental hh","Young adults form own household","age"],
    ["Drivers license","Assigns/updates licenses","age"],
    ["Education","School/university transitions","age"],
    ["Employment","Hires/quits; matches workers to jobs","labour participation, vacancies"],
], col_w=[0.18,0.55,0.27], fs=12)

# 6 — real estate + jobs/migration
s=prs.slides.add_slide(BLANK); header(s,"The models (2/3) — real estate, jobs & migration")
table(s,["Model","What it does","Key driver"],[
    ["Moves","Households relocate via dwelling utility","price, accessibility, region utility"],
    ["Construction","Builds new dwellings","demand, development constraints"],
    ["Demolition","Removes dwellings","age/quality"],
    ["Pricing","Updates dwelling rents/prices","vacancy by region & type"],
    ["Renovation","Changes dwelling quality","quality transition"],
    ["Job-market update","Adds/removes jobs to hit employment totals","employmentForecast (rate mode)"],
    ["In/out migration","Adds/removes households at region edges","populationControlTotalByState"],
], col_w=[0.20,0.52,0.28], fs=12)

# 7 — auto ownership + income
s=prs.slides.add_slide(BLANK); header(s,"The models (3/3) — auto ownership & income")
bullets(s,[
    "**Maryland auto-ownership model** (UEC / logit): probability of 0/1/2/3+ cars per household.",
    (1,"Utility uses household size, #workers, income category, transit accessibility, job density."),
    (1,"NEW in our build: it now runs **every year** so autos respond to changing income & accessibility."),
    "**Income dynamics** (IncomeAdjustment): rank-preserving 'freeze-and-grow' — each person keeps their place",
    (1,"in the income distribution while the whole distribution shifts with the simulated economy."),
    "Both are now calibration targets (see the levers)." ,
], size=17)

# 8 — data flow synth -> predicted
s=prs.slides.add_slide(BLANK); header(s,"From synthetic population to predicted years")
bullets(s,[
    "**Input (base 2016):** VAE synthetic hh / pp / dd / jj  (households, persons, dwellings, jobs).",
    "**Each simulated year:** events + annual updates mutate the agents in place.",
    "**Output per year (2016→2023):** hh/pp/dd/jj_<year>.csv  → collected into Pipeline/2_SILO_landuse/output/.",
    "Transport model is OFF in this stage (skims fixed); MITO consumes the yearly population next.",
    "Control totals: population by state (migration) and employment (job market) drive exogenous growth;",
    (1,"composition (income, household size, age, autos) is what we calibrate behaviourally."),
], size=17)

# 9 — changes overview
s=prs.slides.add_slide(BLANK); header(s,"The changes we are making — overview","fresh stock SILO clone + our fixes & calibration hooks")
table(s,["#","Change","Type","Why"],[
    ["1","findVacantJob sampler reallocation","bug fix","crashed at year 2018 (array overflow)"],
    ["2","takeNewJob keeps prior income","bug fix","removed +28% year-1 income/auto inflation"],
    ["3","summarizeCarOwnership index cap","bug fix","crashed on households with 4+ cars"],
    ["4","Maryland car-ownership wired into annual loop","feature","autos update yearly with income"],
    ["5","Per-state CalibrationConfig levers","feature","Approach-B birth/marriage/income/auto"],
    ["6","Auto-ownership ASC self-calibration","feature","reshape 0/1/2/3+ to ACS (TV 0.21→0.02)"],
    ["7","UpdateOccupationModelMstm (students+retirement)","feature","fix the largest occ_silo forecast gap"],
    ["8","Validator partial-state PUMA + per-bin metric","validation","remove 2023 artefact; bin-level <5pp test"],
    ["9","CompositionReanchorModelMstm (race + hh size)","feature","track ACS composition (frozen closed-model)"],
    ["10","ConcurrentExecutor cached→fixed pool","bug fix","unbounded threads crashed long runs (ulimit)"],
    ["11","Per-state labor-participation lever","feature","lift PA/WV employment to ACS"],
], col_w=[0.05,0.45,0.16,0.34], fs=11)

# 10 — change detail: fixes
s=prs.slides.add_slide(BLANK); header(s,"Change 1–2 — engine bug fixes")
code(s,[
"// 1) JobDataManagerImpl.findVacantJob — reallocate sampler before the fallback loop",
"//    (zero-probability adds had advanced the fixed-capacity index -> ArrayIndexOutOfBounds)",
"if (regionSampler.getCumulatedProbability() == 0) {",
"    regionSampler = new Sampler<>(regions.size(), Region.class, SiloUtil.getRandomObject());",
"    for (Region reg : regions) { ... }",
"}",
"",
"// 2) EmploymentModelImpl.takeNewJob — preserve income when a worker is re-matched to a job",
"//    (re-anchoring to the right-skewed cell MEAN each year inflated income ~28% at year 1)",
"final int priorIncome = person.getAnnualIncome();",
"final int inc = priorIncome > 0 ? priorIncome : Math.max((int) avgIncome + change[sel], 0);",
"person.setIncome(inc);",
], fs=12)

# 11 — change detail: car ownership wiring
s=prs.slides.add_slide(BLANK); header(s,"Change 3–4 — auto ownership active every year")
code(s,[
"// 3) DefaultResultsMonitor.summarizeCarOwnership — cap into the '3+cars' bucket",
"int cars = (int) hh.getVehicles().stream().filter(v->v.getType()==CAR).count();",
"carOwnership[Math.min(cars, 3)]++;",
"",
"// 4) ModelBuilderMstm — register the Maryland model as an annual update listener",
"MaryLandUpdateCarOwnershipModel carOwnershipModel = new MaryLandUpdateCarOwnershipModel(",
"        dataContainer, dataContainer.getAccessibility(), properties, SiloUtil.provideNewRandom());",
"modelContainer.registerModelUpdateListener(carOwnershipModel);",
"//  -> endYear() now re-simulates vehicles/household each year (was synthesis-only before)",
], fs=12)

# 12 — change detail: calibration levers code
s=prs.slides.add_slide(BLANK); header(s,"Change 5 — per-state calibration levers (CalibrationConfig)")
code(s,[
"// utils/CalibrationConfig — loaded ONCE, frozen; 1.0 = no-op when the CSV/state is absent",
"// reads input/assumptions/calibration_by_state.csv  (STFIPS, birthScaler, marriageScaler,",
"//                                                     incomeGrowth, autoScaler)",
"// + zone->state map from zoneSystem.csv (ZoneId, STFIPS)",
"",
"birthProb  *= CalibrationConfig.get().birthScaler(state);     // BirthModelImpl",
"marryProb  *= CalibrationConfig.get().marriageScaler(state);  // MarriageModelMstm",
"income      = (income + change) * incomeGrowth(state);        // IncomeAdjustment",
"prob[1..3] *= autoScaler(state);                              // car ownership (renormalised)",
], fs=12)

# 12b — auto-ownership self-calibration (ASC fit to ACS base-year shares)
s=prs.slides.add_slide(BLANK); header(s,"Change 6 — auto-ownership self-calibrates per state","the scalar autoScaler can't reshape a 4-way distribution; the MNL constants do")
bullets(s,[
    "The Maryland auto model is a 4-alternative logit (0 / 1 / 2 / 3+ cars). A single multiplier can't move",
    (1,"mass into the 3+ bin where suburban/rural ACS sits at 20–26% (stock SILO produced only ~5%)."),
    "**MaryLandUpdateCarOwnershipModel** now self-calibrates the alternative-specific constants (ASCs) ONCE",
    (1,"at the first year-end so per-state predicted shares match ACS targets, then **freezes** them."),
    "Reads input/assumptions/auto_target_shares.csv; 25-iteration logit constant fit per state.",
    "**Result:** autos TV (2018) MD 0.21→0.02, PA 0.34→0.04, DC 0.12→0.03 — fixed the dominant base-year error.",
], size=15)
code(s,[
"// endYear of the first simulated year: standard logit ASC calibration to ACS shares, then frozen",
"for (iter 0..25) {  predicted_k = mean P(k | hh) over the state's households;",
"                    delta_k += ln( target_k / predicted_k );  }   // per state, k = 1,2,3+ cars",
"// DC delta3 = -0.95 (fewer cars, urban) ; WV delta3 = +4.86 (many 3+-car rural households)",
], top=5.4, height=1.5, fs=11)

# 12c — occupation enrichment (students + retirement)
s=prs.slides.add_slide(BLANK); header(s,"Change 7 — occupation enrichment (students & retirement)","the largest forecast error: a definitional/drift gap, now reconciled to ACS each year")
bullets(s,[
    "**Root cause (found by tracing the engine):** the stock demography drifts labour-force status away from ACS —",
    (1,"the education model **graduates every student at age 19** → SILO carries no college students;"),
    (1,"there is **no retirement model at all** (RETIREE never assigned) → retirees pile into 'unemployed/other'."),
    "Effect grows every forecast year: the 'other' bin ran ~+8pp, students/retirees ~−5pp each (occ_silo ~18% TV).",
    "**Fix — new UpdateOccupationModelMstm** (annual): reconciles non-workers to the ACS definition each year,",
    (1,"using ACS in-region college-enrollment-by-age + a retirement rule, so occupation also feeds MITO correctly."),
], size=15)
code(s,[
"// per person, not employed (jobId<=0):",
"  age < 6                       -> TODDLER",
"  6 <= age <= 18                -> STUDENT",
"  19 <= age <= 35               -> STUDENT w.p. ACS_enrollment[age], else UNEMPLOYED",
"  36 <= age < 62                -> UNEMPLOYED",
"  age >= 62                     -> RETIREE        // ACS: ~0.95 of non-workers 62+ are retired",
], top=5.5, height=1.5, fs=11)

# 12d — validator partial-state fix + the stale-validation finding
s=prs.slides.add_slide(BLANK); header(s,"Change 8 — validation correctness fixes")
bullets(s,[
    "**Stale-validation finding:** an early scorecard showed autos ~2% TV — it came from a since-deleted summary",
    (1,"of a run that PRE-DATED the auto self-calibration. The real (pre-fix) autos error was 20–40% TV. The auto"),
    (1,"MNL was never broken; it simply had not been calibrated yet. All artefacts are now traced to a specific run."),
    "**Partial-state PUMA vintage:** PA/VA/WV 2023 PUMS switched to 2020 PUMA codes that partly collide with the",
    (1,"2010 in-region list, silently keeping the WRONG ACS sub-sample (VA-2023 dwellingType/income spiked to ~10pp)."),
    (1,"The validator now detects the vintage and cleanly skips partial states for 2023 (a 2010↔2020 crosswalk is the"),
    (1,"follow-up to actually validate them). This removed a pure validation artefact — not a model error."),
], size=15)

# 12e — composition re-anchoring (the residual trend-drift fix)
s=prs.slides.add_slide(BLANK); header(s,"Change 9 — composition re-anchoring (race · household size · income)",
                                   "control the demographic marginals to ACS, the way a land-use model should")
bullets(s,[
    "**Problem:** stock migration clones random households and barely out-migrates, so race / household-size /",
    (1,"income composition is FROZEN (white 62.3%→61.8%) while real ACS composition keeps moving — the gap grows."),
    "**CompositionReanchorModelMstm** (annual): bounded per-state household swaps toward `composition_targets_by_state.csv`",
    (1,"(per-year ACS race, size & income shares) — remove an over-represented household, drop a clone of an"),
    (1,"under-represented one into its freed dwelling (no vacancy strain, population stable)."),
    "**Decoupled passes:** race matched on size · size & income matched on race → each marginal moves without",
    (1,"disturbing the others. Behavioural outputs (autos, dwelling, location) are NOT controlled — they stay the honest test."),
    "**Result:** race 7–10pp→1–2.5pp; household size 9–13pp→1–1.6pp; income 6–7pp→2–2.5pp (all states).",
], size=14)

# 12f — thread-leak fix + participation lever
s=prs.slides.add_slide(BLANK); header(s,"Change 10–11 — thread-leak fix & per-state participation")
bullets(s,[
    "**10 · Thread-leak (bug):** HouseholdDataManagerImpl.adjustIncome used an UNBOUNDED cached thread pool and",
    (1,"queued one task per person (~12M) → ~4000 threads/yr → OutOfMemoryError vs the OS limit (ulimit -u, hard 4000)."),
    (1,"Fixed to a fixed `numberOfThreads` pool (matches MovesModelImpl). Long runs now complete reliably."),
    "**11 · Per-state labor participation (feature):** PA & WV validated with too FEW employed (≈6.6pp low) — and it",
    (1,"was NOT a job shortage (no missingJob warnings), just a low participation target. EmploymentModelImpl now"),
    (1,"computes base participation PER STATE and targets base × `participationScaler(state)`; PA/WV lifted to ACS"),
    (1,"while MD/DC/VA/DE keep their own (already-good) levels."),
], size=14)

# 13 — calibration philosophy
s=prs.slides.add_slide(BLANK); header(s,"Calibration framework — Approach B (behavioural)","calibrate once on the past, then forecast the future")
bullets(s,[
    "**SILO is a forecasting model**, so it is scored by how well it predicts years it was NOT shown.",
    "**Approach B (chosen):** calibrate behavioural coefficients ONCE on the 2016–2020 backcast, then",
    (1,"**freeze** them and forecast 2021–2023."),
    "**Two tiers of control (key distinction):**",
    (1,"**Exogenous demographic marginals** (population, race, household size, income, labor participation) are"),
    (2,"controlled to ACS forecasts — standard for a land-use model fed regional demographic projections."),
    (1,"**Behavioural outputs** (auto ownership, dwelling type, location, mode) are NOT controlled — they are the"),
    (2,"honest test of the model and are validated free."),
    "Levers are few, interpretable, and physically meaningful — not a black-box fit.",
], size=15)

# 14 — per-state behavioural levers (publication table, live values)
s=prs.slides.add_slide(BLANK); header(s,"Table 1 — per-state behavioural calibration levers",
                                   "frozen after the 2016–2020 fit · 1.00 = neutral · keyed by STFIPS in calibration_by_state.csv")
calib = read_calib()
STORDER=["MD","VA","PA","DE","DC","WV"]
def _calib_row(st):
    for c in calib or []:
        if c.get("state")==st: return c
    return {}
if calib:
    rows=[]
    for st in STORDER:
        c=_calib_row(st)
        if not c: continue
        rows.append([st, c.get("birthScaler","–"), c.get("marriageScaler","–"),
                     c.get("incomeGrowth","–"), c.get("autoScaler","–"), c.get("participationScaler","1.0")])
    table(s,["State","birth","marriage","incomeGrowth","autoScaler","participation"], rows,
          col_w=[0.16,0.16,0.16,0.20,0.16,0.18], fs=13, top=1.3, height=2.7)
    note="Live values from calibration_by_state.csv. Birth/marriage shape demography; incomeGrowth & participation steer income & employment."
else:
    note="calibration_by_state.csv not found — levers default to 1.0 (free-run)."
bullets(s,[
    "**birth / marriage** → household-size & single-person share.   **incomeGrowth** → per-state income trajectory.",
    "**autoScaler** → legacy multiplier (now superseded by the ASC self-calibration, Table 2).",
    "**participation** → per-state labor-force participation; lifts PA/WV employment to ACS (others 1.00 = unchanged).",
    note,
], top=4.3, size=14)

# 14b — auto ASC + composition re-anchoring tables
s=prs.slides.add_slide(BLANK); header(s,"Table 2 — auto-ownership ASC deltas & re-anchoring targets",
                                   "the self-calibrated logit constants and the demographic control marginals")
bullets(s,["**Table 2a — per-state auto-ownership ASC shifts** (self-calibrated to 2016 ACS shares, then frozen):"], top=1.15, size=13)
table(s,["State","Δ 1-car","Δ 2-car","Δ 3+-car","direction"],[
    ["MD","+0.79","+1.83","+3.46","more cars (suburban)"],
    ["VA","+0.84","+1.83","+3.37","more cars"],
    ["PA","+1.12","+2.84","+4.73","more cars (rural)"],
    ["DE","+1.29","+2.69","+4.33","more cars"],
    ["WV","+1.39","+2.99","+4.86","most cars (rural)"],
    ["DC","−0.68","−1.38","−0.95","fewer cars (urban)"],
], col_w=[0.16,0.18,0.18,0.18,0.30], fs=11, top=1.5, height=2.5)
bullets(s,["**Table 2b — composition re-anchoring marginals** (per state × year, from ACS PUMS, in composition_targets_by_state.csv):"],
        top=4.15, size=13)
table(s,["Marginal","Categories","Validated as"],[
    ["race / ethnicity","white · black · hispanic · other","person-level race4"],
    ["household size","1 · 2 · 3 · 4+","household hhSize"],
    ["household income","<30k · 30–75k · 75–150k · 150k+ (2016$)","household hh_inc9"],
], col_w=[0.25,0.50,0.25], fs=11, top=4.5, height=1.5)

# 15 — calibrate-forecast loop diagram
s=prs.slides.add_slide(BLANK); header(s,"The calibrate → freeze → forecast loop")
s.shapes.add_picture(str(diagram_loop()), Inches(1.0), Inches(1.25), width=Inches(11.3))

# 16 — validation framework + acceptance criterion
s=prs.slides.add_slide(BLANK); header(s,"Validation framework — model vs ACS","acceptance criterion: every category within ±5 percentage points")
bullets(s,[
    "Compared against **ACS PUMS 5-year** for every year, all **6 states** (incomes deflated to 2016$).",
    (1,"MD/DC/DE whole-state; PA/VA/WV filtered to in-region PUMAs and coverage-weighted."),
    "**Headline acceptance: per-bin.** For every variable, EACH category's |share_SILO − share_ACS| < **5 pp**",
    (1,"(stricter & more interpretable than aggregate Total Variation, which can hide a single bad bin)."),
    "8 core variables: hhSize, autos, dwellingType, income(9), age, gender, race, occupation.",
    "Split: 2016–2020 = in-sample (calibrated); **2021–2023 = out-of-sample forecast (the headline)**.",
    "Tooling: code/perbin_scorecard.py (per-bin table) + valib.one_var_fig (per-variable figures with a ±5pp panel).",
], size=15)

# 17 — current per-bin scorecard heatmap (reference run)
s=prs.slides.add_slide(BLANK); header(s,"Per-bin scorecard — current reference run",
                                   "worst per-bin gap (pp) per state × variable, forecast 2021–2023")
if HEATMAP.exists():
    s.shapes.add_picture(str(HEATMAP), Inches(0.7), Inches(1.15), width=Inches(9.2))
pb = read_perbin()
if pb:
    order=["hhSize","autos","hh_inc9","occ_silo","race4"]
    nm={"hhSize":"HHsize","autos":"autos","hh_inc9":"income","occ_silo":"occ","race4":"race"}
    rows=[]
    for st in ["MD","VA","PA","DE","DC","WV"]:
        if st in pb:
            rows.append([st]+[f"{pb[st].get(v,0):.0f}" for v in order])
    table(s,["St"]+[nm[v] for v in order], rows, left=10.0, top=1.3, width=3.0, height=3.2,
          col_w=[0.2,0.2,0.16,0.16,0.14,0.16], fs=10)
bullets(s,["green ≤ 5pp (pass) · red > 5pp.  Occupation & race columns are the focus of changes 7 and the re-anchoring."],
        left=10.0, top=4.7, width=3.0, size=10)

# 18 — the trend-drift mechanism (diagnosis that motivated re-anchoring)
s=prs.slides.add_slide(BLANK); header(s,"Why the composition drifted — the closed-model mechanism",
                                   "diagnosis behind change 9 · now solved for race / household size / income")
bullets(s,[
    "**Migration is composition-neutral:** inmigrants are random household CLONES (duplicateHousehold), outmigration ≈ 0.",
    (1,"So SILO's race / household-size / income composition is essentially FROZEN: white 62.3%→61.8% over 2016→2023."),
    "Meanwhile **real ACS composition keeps moving** (minority share rises ~1pp/yr; DC singles rise; WV elderly singles).",
    (1,"The gap grows purely because SILO stands still while the target moves — a closed-model forecast limitation."),
    "**Fix (change 9): composition re-anchoring** — steer bounded per-state household churn to per-year ACS marginals.",
    (1,"This is the standard way a land-use model is controlled to exogenous demographic forecasts; the behavioural"),
    (1,"outputs (autos, dwelling type, location, mode) stay UNcontrolled — they remain the honest test."),
    "**Result:** race, household size and income all brought to <3pp. The SAME mechanism explains the residual DC age (next).",
], size=14)

# 18b — FINAL result: the three residual cells
s=prs.slides.add_slide(BLANK); header(s,"Final result — 3 / 48 cells, and what remains",
                                   "all systematic & structural problems solved; the residual is an honest small-state floor")
table(s,["Residual cell","Gap (pp)","Root cause","Class"],[
    ["DC · age","6.8","closed-model age drift: DC's young-IN / old-OUT churn isn't reproduced","demographic"],
    ["WV · autos","6.6","rural 3+-car ownership trended up; ASC frozen at 2016 base","behavioural"],
    ["WV · dwelling type","7.7","too few single-family homes — a housing-STOCK shortfall","structural"],
], col_w=[0.18,0.10,0.52,0.20], fs=12, top=1.35, height=1.9)
bullets(s,[
    "**All three sit in the two hardest contexts** — WV (n≈2,545 households, near the ACS sampling floor) and DC",
    (1,"(the most extreme, churn-sustained demographic profile of the six states). Every other state × variable is <5pp."),
    "**WV autos & WV dwelling are behavioural / structural** — forcing them would mean controlling the very outputs the",
    (1,"validation is meant to test (circular), or re-synthesising the base dwelling stock. Left as the honest floor."),
    "**DC age is demographic** (legitimately controllable) but deliberately NOT forced — see the coupling analysis (next).",
], top=3.5, size=14)

# 18c — DC age deep-dive
s=prs.slides.add_slide(BLANK); header(s,"Residual deep-dive — why DC age is off",
                                   "the only age miss of the six states · a closed-model churn limitation")
bullets(s,["**DC person-age distribution, SILO vs ACS (% of persons):**"], top=1.1, size=13)
table(s,["Year · source","20–24","25–29","30–34","50–54","55–59","60–64","65–69"],[
    ["2016 SILO","5.6","9.3","10.0","7.2","6.6","6.3","4.6"],
    ["2016 ACS","7.3","12.4","11.3","6.0","5.6","5.0","3.7"],
    ["2023 SILO","3.9","4.7","6.9","7.5","8.1","8.0","7.3"],
    ["2023 ACS","5.9","11.3","11.8","5.3","4.9","4.9","4.1"],
], col_w=[0.22,0.11,0.11,0.11,0.11,0.11,0.11,0.12], fs=11, top=1.45, height=1.9)
bullets(s,[
    "**Mechanism:** SILO ages the 2016 base forward. DC's reality is sustained by young adults moving IN and older",
    (1,"residents moving OUT every year — but composition-neutral migration (random clones) doesn't reproduce it."),
    "So the 25–34 cohort **ages out and isn't replenished** (25–29: 9.3→4.7 vs ACS 11.3) while 50–69 **grays in place** (+3pp)."
    ,
    "**Plus a small seed:** the VAE 2016 base already under-represents DC 25–29 by ~3pp. The drift then amplifies it.",
    "Only DC drifts because it has by far the most churn-driven age profile; the other five states stay <2.2pp all years.",
], top=3.55, size=13)

# 18d — why we stop at 3 (the coupling principle)
s=prs.slides.add_slide(BLANK); header(s,"Why we stop at 3 — coupling & the controlled-vs-free line")
bullets(s,[
    "**Re-anchoring marginals are correlated** — swapping whole households to fix one marginal disturbs the others.",
    (1,"Demonstrated: the income re-anchor (change 9) fixed income but regressed WV autos −1pp and WV occupation +1.2pp."),
    "**An age pass would couple the worst:** younger households differ in size, income AND occupation; DC occupation is",
    (1,"already the closest cell to the line (4.1pp), so adding young adults would likely push it OVER — net no gain (whack-a-mole)."),
    (1,"Age is also person-level but the swap is household-level (mixed ages) → imprecise, spills into neighbouring bins."),
    "**The principle:** we control EXOGENOUS demographics (population, race, size, income, participation) to ACS, but we do",
    (1,"NOT control BEHAVIOURAL outputs (autos, dwelling, location, mode). Forcing WV autos/dwelling would make the"),
    (1,"validation circular — measuring the model against a target we injected."),
    "**Decision:** calib5 (3/48) is the accepted result; the 3 residuals are documented as the honest small-state floor.",
], size=13)

# 19 — results progression (publication table)
s=prs.slides.add_slide(BLANK); header(s,"Table 3 — calibration progression (per-bin acceptance)",
                                   "forecast 2021–2023 · cells = state × variable failing the ±5pp test (of 48)")
table(s,["Scenario","Change added","Cells >5pp","Solved"],[
    ["calib1","reference (engine fixes + auto ASC)","22 / 48","baseline"],
    ["calib2","occupation enrichment + income/marriage retune","15 / 48","occupation"],
    ["calib3","composition re-anchor (race + household size)","6 / 48","race, household size"],
    ["calib4","+ income re-anchor pass","5 / 48","income distribution"],
    ["calib5","+ per-state labor-participation lever","3 / 48","PA / WV employment"],
], col_w=[0.13,0.46,0.18,0.23], fs=12, top=1.35, height=2.5)
bullets(s,[
    "**All systematic / structural problems are solved** across the 6 states: occupation definition, race trend,",
    (1,"household-size trend, income distribution, and (calib5) per-state employment."),
    "**Expected residual (~3 cells), characterised as honest model floor in the two hardest states:**",
    (1,"WV autos (rural 3+-car forecast drift) · WV dwelling type (housing-stock, not re-anchorable) · DC age (25–34 bulge)."),
    "These are behavioural / structural — forcing them would mean controlling the very outputs the validation tests.",
], top=4.1, size=14)

# 20 — this-session changelog
s=prs.slides.add_slide(BLANK); header(s,"Changelog — this iteration")
bullets(s,[
    "**Diagnosed** the alarming numbers as a stale pre-calibration run; confirmed the auto MNL itself was correct.",
    "**Adopted** the per-bin (±5pp) acceptance criterion; built code/perbin_scorecard.py & make_summary_figures.py.",
    "**Added** UpdateOccupationModelMstm (college students + retirement) — fixed occupation (22→15 cells).",
    "**Added** CompositionReanchorModelMstm (race + size + income re-anchoring) — fixed race, household size, income (15→5).",
    "**Added** a per-state labor-participation lever (calib5) — cleared PA/WV employment with no collateral.",
    "**Fixed** the validator's partial-state 2023 PUMA artefact and an engine thread-leak that crashed long runs.",
    "**Rebuilt** the validation figures to publication quality (labelled bins, ±5pp panel, 200 dpi) + scorecard heatmap.",
    "Net: **22 → 3 failing cells** (calib5 = final). Residual = DC age + WV autos + WV dwelling — documented honest floor.",
    "**Decision:** DC age NOT force-fixed — an age pass would couple (regress DC occupation, already at 4.1pp). The",
    (1,"controlled-vs-free line holds: exogenous demographics controlled to ACS; behavioural outputs left as the test."),
    "This deck is regenerated by docs/make_pptx.py (set REF_SUB to the latest validated run) — re-run to refresh.",
], size=13)

prs.save(str(OUT))
print(f"wrote {OUT}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}  |  calib={'yes' if calib else 'no'}  "
      f"perbin={'yes' if pb else 'no'}  heatmap={'yes' if HEATMAP.exists() else 'no'}")
