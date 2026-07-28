#!/usr/bin/env python3
"""Assemble the v7 base AADT (Simulated vs Observed) figures into ONE 16:9 deck.

Order (per request): the OVERALL pooled figure FIRST, then the hierarchy summaries,
the per-facility scatter/GEH/bias set, the per-route summary table, and finally the
individual per-facility / per-speed-tier / per-route panels arranged in tidy grids.

Reproducible: reads the PNGs already rendered into network_validation_2023/v7_base/
by make_aadt_route_figures_v7.py + run_v7_base_validation.py. Rerun after re-rendering.
Writes: network_validation_2023/v7_base/AADT_Simulated_vs_Observed_v7.pptx
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path("/Users/tomal/Documents/SILO MITO Chayan/VAE-SILO-MITO-MATSIM/Updated MATSim")
V7   = ROOT/"network_validation_2023/v7_base"
OUT  = V7/"AADT_Simulated_vs_Observed_v7.pptx"

EMU_IN = 914400
NAVY   = RGBColor(0x1F, 0x38, 0x64)
GREY   = RGBColor(0x55, 0x55, 0x55)
DARK   = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def _txt(slide, x, y, w, h, text, size, *, bold=False, color=DARK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Calibri"
    return tb

def _fit(img_path, cx, cy, cw, ch):
    """Return (x,y,w,h) EMU fitting the image into the (cx,cy,cw,ch) box, centered, aspect-kept."""
    iw, ih = Image.open(img_path).size
    box_ar = cw/ch; img_ar = iw/ih
    if img_ar >= box_ar:
        w = cw; h = int(cw/img_ar)
    else:
        h = ch; w = int(ch*img_ar)
    x = cx + (cw-w)//2; y = cy + (ch-h)//2
    return int(x), int(y), int(w), int(h)

def pic(slide, path, cx, cy, cw, ch):
    x, y, w, h = _fit(str(path), cx, cy, cw, ch)
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)

def _band(slide, y, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SW, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def title_slide():
    s = prs.slides.add_slide(BLANK)
    _band(s, Inches(2.55), Inches(2.4), NAVY)
    _txt(s, Inches(0.8), Inches(2.75), Inches(11.7), Inches(1.0),
         "Simulated vs Observed Daily Traffic (AADT)", 34, bold=True,
         color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    _txt(s, Inches(0.8), Inches(3.65), Inches(11.7), Inches(0.7),
         "Baltimore Region, 2023  ·  MATSim v7 base run (resident-only demand, 10% sample scaled ×10)",
         16, color=RGBColor(0xDD,0xE4,0xF0), align=PP_ALIGN.CENTER)
    _txt(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.2),
         "Observed = passenger-car AADT (freight/bus removed via MDOT 2023 vehicle-class shares where "
         "available; facility-median fallback elsewhere).  NCHRP 255/765 facility bands and GEH shown for "
         "reference. Interstates under-count is consistent with through / non-resident passenger traffic "
         "outside the resident-only scope (hypothesis).", 12, color=GREY, align=PP_ALIGN.CENTER)

def bullets_slide(title, bullets, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    _txt(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.75), title, 26, bold=True, color=NAVY)
    if subtitle:
        _txt(s, Inches(0.62), Inches(1.12), Inches(12.1), Inches(0.5), subtitle, 14, italic=True, color=GREY)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(11.9), Inches(5.3))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lead, body) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r = p.add_run(); r.text = "•  " + lead
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"
        if body:
            r2 = p.add_run(); r2.text = "  " + body
            r2.font.size = Pt(14); r2.font.bold = False; r2.font.color.rgb = DARK; r2.font.name = "Calibri"
    return s

def section_slide(title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    _txt(s, Inches(0.7), Inches(3.0), Inches(12.0), Inches(1.0), title, 30, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)
    if subtitle:
        _txt(s, Inches(0.7), Inches(4.05), Inches(12.0), Inches(0.8), subtitle, 15,
             color=GREY, align=PP_ALIGN.CENTER)
    ln = _band(s, Inches(3.9), Pt(2.5), NAVY)

def one_fig(title, path, note=""):
    s = prs.slides.add_slide(BLANK)
    _txt(s, Inches(0.45), Inches(0.18), Inches(12.45), Inches(0.6), title, 22, bold=True, color=NAVY)
    top = Inches(0.95); bot_note = Inches(0.5) if note else Inches(0.15)
    box_h = SH - top - bot_note
    pic(s, path, Inches(0.35), top, SW-Inches(0.7), box_h)
    if note:
        _txt(s, Inches(0.45), SH-Inches(0.5), Inches(12.45), Inches(0.4), note, 11,
             italic=True, color=GREY, align=PP_ALIGN.CENTER)

def grid_figs(title, items, ncol, nrow):
    """items: list of (caption, path). Lays out ncol x nrow per slide, paginating."""
    per = ncol*nrow
    for pg in range(0, len(items), per):
        chunk = items[pg:pg+per]
        s = prs.slides.add_slide(BLANK)
        ttl = title + (f"  ({pg//per+1}/{(len(items)+per-1)//per})" if len(items) > per else "")
        _txt(s, Inches(0.45), Inches(0.16), Inches(12.45), Inches(0.55), ttl, 20, bold=True, color=NAVY)
        gx, gy = Inches(0.3), Inches(0.85)
        gw, gh = SW-Inches(0.6), SH-Inches(1.05)
        cap_h = Inches(0.32)
        cw = gw/ncol; ch = gh/nrow
        for i, (cap, path) in enumerate(chunk):
            r, c = divmod(i, ncol)
            cx = int(gx + c*cw); cy = int(gy + r*ch)
            if cap:
                _txt(s, cx, cy, int(cw), cap_h, cap, 11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
                pic(s, path, cx+Inches(0.05), cy+cap_h, int(cw)-Inches(0.1), int(ch)-cap_h-Inches(0.05))
            else:
                pic(s, path, cx+Inches(0.05), cy, int(cw)-Inches(0.1), int(ch)-Inches(0.05))

# ---------------------------------------------------------------- BUILD
title_slide()

# 1) OVERALL figure FIRST
one_fig("Overall — All count stations pooled", V7/"all_stations_sim_vs_obs.png",
        "Every non-ramp station, all facility classes pooled · single ±50% acceptance band (0.5×–1.5×).")

# 1b) Regional inflow/outflow — gateway station selection (considerations) + map
try:
    import pandas as _pd
    _sc = _pd.read_csv(V7/"screenline_hybrid.csv")
    _tot = _sc[_sc.station.astype(str).str.contains("TOTAL")].iloc[0]
    _scr = f"Σsim {_tot.sim:,.0f} vs Σobs {_tot.obs:,.0f}  =  {_tot.diff_pct:+.1f}%"
except Exception:
    _scr = "see screenline_hybrid.csv"
bullets_slide("Regional inflow / outflow — gateway station selection",
    [("Why a gateway screenline.",
      "Demand is RESIDENT-ONLY — external and through trips are not simulated — so absolute link "
      "volumes run low on high-through corridors. Total flow entering/leaving the region is checked "
      "with a RADIAL SCREENLINE of the count stations on the corridors that carry that in/out flow."),
     ("The 14 gateways selected.",
      "I-95 (SW & NE), I-83 (N), I-70 (W), US-40 (W & E), MD-295 (S), I-795 (NW), I-97 (S), "
      "MD-140 (NW), MD-26 (W), MD-2 (S), MD-43 (NE), MD-144 (W) — one clean count per principal radial approach."),
     ("Selection criteria.",
      "(a) a limited-access or major-arterial RADIAL carrying regional in/out flow; (b) a clean station→link "
      "match (no ramp / opposite-carriageway / over-tolerance snap); (c) spread across compass approaches so "
      "every gate direction is represented."),
     ("Scope caveat.",
      "These are RADIAL crossings, NOT a closed beltway cordon — the screenline Σ is a directional in/out "
      "consistency check, not a mass-balance cordon. One near-total miss (B0988, I-95 NE) is flagged and drags the sum."),
     ("Radial screenline result.", f"{_scr}  (resident-only; the ~−39% deficit matches the resident-scope band).")],
    subtitle="The stations used to push / check regional inflow and outflow (validate_base_hybrid NAMED + EXTRA).")
one_fig("Radial inflow/outflow gateway stations — regional map", V7/"gateway_stations_map.png",
        "All AADT stations coloured by facility; the 14 radial inflow/outflow gateways marked as red stars (route + approach direction).")

# 2) Hierarchy summaries
one_fig("AADT across the FHWA facility hierarchy", V7/"facility_hierarchy_summary.png")
one_fig("AADT by design (free-flow) speed tier", V7/"speedtier_hierarchy_summary.png")

# 3) By-facility scatter / GEH / bias set
one_fig("Scatter by facility class (total AADT)", V7/"figA_scatter_by_facility.png")
grid_figs("GEH<5 share and volume bias by facility class",
          [("GEH<5 share", V7/"figB_geh_by_facility.png"),
           ("Median relative bias", V7/"figC_relbias_by_facility.png")], 2, 1)

# 4) Per-route summary table
one_fig("Per-route metrics summary", V7/"aadt_validation_by_route/route_validation_summary_table.png")

# 5) Individual panels
section_slide("Facility-tier panels", "FHWA F_SYSTEM hierarchy (ramps excluded)")
fac_order = ["Interstate","OtherFreewayExpressway","PrincipalArterial","MinorArterial",
             "MajorCollector","MinorCollectorLocal"]
fac_names = {"Interstate":"Interstate","OtherFreewayExpressway":"Other Freeway-Expressway",
             "PrincipalArterial":"Principal Arterial","MinorArterial":"Minor Arterial",
             "MajorCollector":"Major Collector","MinorCollectorLocal":"Minor Collector-Local"}
grid_figs("Facility-tier panels",
          [(fac_names[k], V7/f"by_facility/facility_{k}.png") for k in fac_order], 3, 2)

section_slide("Design-speed-tier panels", "Stations binned by matched-link free-flow speed")
spd_order = ["Freeway","MajorArterial","Arterial","Collector","LocalStreet"]
spd_names = {"Freeway":"Freeway (≥55 mph)","MajorArterial":"Major Arterial (45–55)",
             "Arterial":"Arterial (35–45)","Collector":"Collector (25–35)","LocalStreet":"Local Street (<25)"}
grid_figs("Design-speed-tier panels",
          [(spd_names[k], V7/f"by_speedtier/speedtier_{k}.png") for k in spd_order], 3, 2)

section_slide("Per-route panels", "Named interstates, US routes, and major MD arterials")
RD = V7/"aadt_validation_by_route"
route_order = ["I95","I695","I83","I70","I895","I97","I795","MD295","US1","US40",
               "MD2","MD140","MD45","MD144","MD26","MD170","MD139","MD25","MD648","MD3",
               "MD175","MD97","MinorArterial(all)","CollectorLocal(all)"]
route_names = {"I95":"I-95","I695":"I-695","I83":"I-83","I70":"I-70","I895":"I-895","I97":"I-97",
               "I795":"I-795","MD295":"MD-295","US1":"US-1","US40":"US-40","MD2":"MD-2","MD140":"MD-140",
               "MD45":"MD-45","MD144":"MD-144","MD26":"MD-26","MD170":"MD-170","MD139":"MD-139",
               "MD25":"MD-25","MD648":"MD-648","MD3":"MD-3","MD175":"MD-175","MD97":"MD-97",
               "MinorArterial(all)":"Minor Arterial (all)","CollectorLocal(all)":"Collector-Local (all)"}
route_items = [(route_names[k], RD/f"aadt_{k}.png") for k in route_order
               if (RD/f"aadt_{k}.png").exists()]
grid_figs("Per-route panels", route_items, 3, 2)

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
