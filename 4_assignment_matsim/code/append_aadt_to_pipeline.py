#!/usr/bin/env python3
"""Update the AADT appendix inside the master pipeline deck.

  VAE-SILO-MITO-MATSim Pipeline.pptx

Removes the OLD AADT figure appendix (the per-route / per-tier slides carrying the
[diagnostic] / [validates] / [reference] tags) and appends the UPDATED section:
  overall pooled figure FIRST -> gateway considerations + network map -> hierarchy
  summaries -> figA/B/C -> route table -> every per-facility / per-speed / per-route
  panel, ONE PICTURE PER SLIDE, with neutral wording.

Idempotent: rerun after re-rendering figures; it strips any prior updated block too
(it removes from the first AADT figure slide to the end before appending).
A timestamped backup of the deck is written first.
"""
import shutil, datetime as dt
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

REPO = Path("/Users/tomal/Documents/SILO MITO Chayan/VAE-SILO-MITO-MATSIM")
DECK = REPO/"VAE-SILO-MITO-MATSim Pipeline.pptx"
V7   = REPO/"Updated MATSim/network_validation_2023/v7_base"
RD   = V7/"aadt_validation_by_route"

NAVY = RGBColor(0x1F,0x38,0x64); GREY = RGBColor(0x55,0x55,0x55); DARK = RGBColor(0x22,0x22,0x22)

# backup
bk = DECK.with_name(f"VAE-SILO-MITO-MATSim Pipeline_backup_{dt.datetime.now():%Y%m%d_%H%M}.pptx")
shutil.copy2(DECK, bk); print("backup ->", bk.name)

prs = Presentation(str(DECK))
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]   # 'Blank'

def _title_of(s):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.split("\n")[0]
    return ""

# --- 1) rename the AADT section lead-in slide (strip "Validation") --------------------------------
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("AADT-2023 Validation"):
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if "AADT-2023 Validation" in r.text:
                        r.text = r.text.replace("AADT-2023 Validation — Honest Resident Scope",
                                                "AADT — Simulated vs Observed (Honest Resident Scope)")

# --- 2) delete the OLD AADT figure appendix (first tagged slide -> end) ----------------------------
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
start = None
for i, s in enumerate(prs.slides):
    t = _title_of(s)
    if any(tag in t for tag in ("[diagnostic]", "[validates]", "[reference]")):
        start = i; break
if start is not None:
    print(f"removing old AADT appendix: slides {start+1}..{len(ids)} ({len(ids)-start} slides)")
    for sldId in ids[start:]:
        prs.part.drop_rel(sldId.get(qn("r:id"))); sldIdLst.remove(sldId)
else:
    print("no tagged old-appendix slides found; appending at end")

# ---------------------------------------------------------------- slide builders
def _txt(slide,x,y,w,h,text,size,*,bold=False,color=DARK,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,italic=False):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
    f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name="Calibri"
    return tb

def _fit(path,cx,cy,cw,ch):
    iw,ih=Image.open(path).size; box=cw/ch; im=iw/ih
    if im>=box: w=cw; h=int(cw/im)
    else: h=ch; w=int(ch*im)
    return int(cx+(cw-w)//2),int(cy+(ch-h)//2),int(w),int(h)

def one_fig(title,path,note=""):
    s=prs.slides.add_slide(BLANK)
    _txt(s,Inches(0.45),Inches(0.16),Inches(12.45),Inches(0.6),title,22,bold=True,color=NAVY)
    top=Inches(0.92); botn=Inches(0.5) if note else Inches(0.12)
    x,y,w,h=_fit(str(path),Inches(0.35),top,SW-Inches(0.7),SH-top-botn)
    s.shapes.add_picture(str(path),x,y,width=w,height=h)
    if note:
        _txt(s,Inches(0.45),SH-Inches(0.48),Inches(12.45),Inches(0.4),note,11,italic=True,
             color=GREY,align=PP_ALIGN.CENTER)

def section(title,subtitle=""):
    s=prs.slides.add_slide(BLANK)
    _txt(s,Inches(0.7),Inches(2.95),Inches(12.0),Inches(1.0),title,30,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
    if subtitle:
        _txt(s,Inches(0.7),Inches(4.0),Inches(12.0),Inches(0.8),subtitle,15,color=GREY,align=PP_ALIGN.CENTER)

def bullets(title,items,subtitle=""):
    s=prs.slides.add_slide(BLANK)
    _txt(s,Inches(0.6),Inches(0.35),Inches(12.1),Inches(0.75),title,26,bold=True,color=NAVY)
    if subtitle: _txt(s,Inches(0.62),Inches(1.12),Inches(12.1),Inches(0.5),subtitle,14,italic=True,color=GREY)
    tb=s.shapes.add_textbox(Inches(0.7),Inches(1.75),Inches(11.9),Inches(5.3)); tf=tb.text_frame; tf.word_wrap=True
    for i,(lead,body) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(9)
        r=p.add_run(); r.text="•  "+lead; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=NAVY; r.font.name="Calibri"
        if body:
            r2=p.add_run(); r2.text="  "+body; r2.font.size=Pt(14); r2.font.color.rgb=DARK; r2.font.name="Calibri"

# ---------------------------------------------------------------- BUILD updated AADT section
section("AADT — Simulated vs Observed (v7 base run, 2023)",
        "Resident-only demand · 10% sample scaled ×10 · passenger-car counts · one figure per slide")

# overall FIRST
one_fig("Overall — all count stations pooled", V7/"all_stations_sim_vs_obs.png",
        "All non-ramp stations pooled · single ±50% band (0.5×–1.5×) · I-695 Beltway stations circled.")

# gateway considerations + network map
try:
    import pandas as _pd
    _sc=_pd.read_csv(V7/"screenline_hybrid.csv"); _t=_sc[_sc.station.astype(str).str.contains("TOTAL")].iloc[0]
    _scr=f"Σsim {_t.sim:,.0f} vs Σobs {_t.obs:,.0f} = {_t.diff_pct:+.1f}%"
except Exception: _scr="see screenline_hybrid.csv"
bullets("Regional inflow / outflow — gateway station selection",
    [("Why a gateway screenline.","Demand is RESIDENT-ONLY (no external/through trips), so absolute link "
      "volumes run low on high-through corridors. Regional inflow/outflow is checked with a RADIAL SCREENLINE "
      "of the count stations on the corridors that carry that in/out flow."),
     ("The 14 gateways.","I-95 (SW & NE), I-83 (N), I-70 (W), US-40 (W & E), MD-295 (S), I-795 (NW), I-97 (S), "
      "MD-140 (NW), MD-26 (W), MD-2 (S), MD-43 (NE), MD-144 (W) — one clean count per principal radial approach."),
     ("Selection criteria.","(a) limited-access/major-arterial RADIAL carrying regional in/out flow; (b) clean "
      "station→link match (no ramp / opposite-carriageway / over-tolerance snap); (c) spread across compass approaches."),
     ("Scope caveat.","RADIAL crossings, NOT a closed beltway cordon — a directional in/out consistency check, not "
      "a mass-balance cordon. One near-total miss (B0988, I-95 NE) is flagged and drags the sum."),
     ("Radial screenline result.",f"{_scr}  (resident-only; the ~−39% deficit matches the resident-scope band).")],
    subtitle="The points selected to check regional inflow / outflow.")
one_fig("Radial inflow/outflow gateways — MATSim network + county map", V7/"gateway_stations_map.png",
        "MATSim road network (freeways emphasized) + county boundaries; 14 radial gateways as red stars (route + approach).")

# hierarchy summaries + facility set + table
one_fig("AADT across the FHWA facility hierarchy", V7/"facility_hierarchy_summary.png")
one_fig("AADT by design (free-flow) speed tier", V7/"speedtier_hierarchy_summary.png")
one_fig("Scatter by facility class (total AADT)", V7/"figA_scatter_by_facility.png")
one_fig("GEH<5 share by facility class", V7/"figB_geh_by_facility.png")
one_fig("Median relative bias by facility class", V7/"figC_relbias_by_facility.png")
one_fig("Per-route metrics summary", RD/"route_validation_summary_table.png")

# facility-tier panels (one each)
section("Facility-tier panels", "FHWA F_SYSTEM hierarchy (ramps excluded)")
fac=[("Interstate","Interstate"),("OtherFreewayExpressway","Other Freeway-Expressway"),
     ("PrincipalArterial","Principal Arterial"),("MinorArterial","Minor Arterial"),
     ("MajorCollector","Major Collector"),("MinorCollectorLocal","Minor Collector-Local")]
for k,nm in fac: one_fig(f"Facility tier — {nm}", V7/f"by_facility/facility_{k}.png")

# speed-tier panels
section("Design-speed-tier panels", "Stations binned by matched-link free-flow speed")
spd=[("Freeway","Freeway (≥55 mph)"),("MajorArterial","Major Arterial (45–55 mph)"),
     ("Arterial","Arterial (35–45 mph)"),("Collector","Collector (25–35 mph)"),
     ("LocalStreet","Local Street (<25 mph)")]
for k,nm in spd: one_fig(f"Design speed — {nm}", V7/f"by_speedtier/speedtier_{k}.png")

# per-route panels
section("Per-route panels", "Named interstates, US routes, and major MD arterials")
routes=[("I95","I-95"),("I695","I-695 (Beltway)"),("I83","I-83"),("I70","I-70"),("I895","I-895"),
        ("I97","I-97"),("I795","I-795"),("MD295","MD-295"),("US1","US-1"),("US40","US-40"),
        ("MD2","MD-2"),("MD140","MD-140"),("MD45","MD-45"),("MD144","MD-144"),("MD26","MD-26"),
        ("MD170","MD-170"),("MD139","MD-139"),("MD25","MD-25"),("MD648","MD-648"),("MD3","MD-3"),
        ("MD175","MD-175"),("MD97","MD-97"),("MinorArterial(all)","Minor Arterial (all)"),
        ("CollectorLocal(all)","Collector-Local (all)")]
for k,nm in routes:
    p=RD/f"aadt_{k}.png"
    if p.exists(): one_fig(f"{nm} — Simulated vs Observed AADT", p)

prs.save(str(DECK))
print(f"saved {DECK.name}  ({len(prs.slides._sldIdLst)} slides total)")
